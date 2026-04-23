
import sys
#import warnings

import anndata
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import scanpy as sc

import scvi
import os
import anndata as ad
import seaborn as sns

import warnings
import copy
import logging
import shutil
import matplotlib.pyplot as plt
#%matplotlib inline
import scipy.sparse as sparse
from joblib import parallel_backend

import scrublet as scr
import subprocess

import SEACells

import matplotlib
import seaborn as sns

import re

# Some plotting aesthetics  ### necessary for plotting within the interactive session.


%matplotlib inline  
sns.set_style('ticks')
matplotlib.rcParams['figure.figsize'] = [8, 8]
matplotlib.rcParams['figure.dpi'] = 200
sc.set_figure_params(figsize=(8, 8))
scvi.settings.seed = 94705
matplotlib.use("module://matplotlib_inline.backend_inline")
%config InlineBackend.print_figure_kwargs={'facecolor' : "w"}
%config InlineBackend.figure_format='retina'



def run_SEACells(ad):
    sc.pp.highly_variable_genes(ad, flavor="seurat_v3",
                            n_top_genes=5000,subset=True)
    scvi.model.SCVI.setup_anndata(ad,continuous_covariate_keys=["n_genes","pct_counts_mt","pct_counts_ribo"])
    model = scvi.model.SCVI(ad)
    model.to_device('cuda:1')   # moves model to GPU 1
    model.train(max_epochs=200,early_stopping=True)
    ad.obsm["X_scVI"] = model.get_latent_representation()
    ## Core parameters 
    n_SEACells = 90
    build_kernel_on = 'X_scVI' # key in ad.obsm to use for computing metacells
                            # This would be replaced by 'X_svd' for ATAC data
    ## Additional parameters
    n_waypoint_eigs = 10 # Number of eigenvalues to consider when initializing metacells
    model = SEACells.core.SEACells(ad, 
                    build_kernel_on=build_kernel_on, 
                    n_SEACells=n_SEACells, 
                    n_waypoint_eigs=n_waypoint_eigs,
                    convergence_epsilon = 1e-5)
    model.construct_kernel_matrix()
    M = model.kernel_matrix
    model.initialize_archetypes()
    model.fit(min_iter=10, max_iter=100)
    with parallel_backend('threading', n_jobs=16):
        sc.pp.neighbors(adata, use_rep="X_scVI")
        sc.tl.leiden(adata,flavor="igraph", n_iterations=-1,directed=False)
        sc.tl.umap(adata)
    return model, ad



import pandas as pd
import scipy.stats as stats

def compute_barcode_overlap_with_pval_and_filter(adata_list, sample_names=None, total_barcodes=1_000_000, pval_threshold=1e-8):
    """
    Compute cell barcode overlap and p-values using a hypergeometric test.
    Select sample pairs where P < 1E-08.
    Parameters:
        adata_list (list of AnnData): List of AnnData objects.
        sample_names (list of str, optional): Sample names corresponding to each AnnData object.
        total_barcodes (int): Total barcode space, assumed to be 1M.
        pval_threshold (float): Threshold for selecting significant sample pairs.

    Returns:
        overlap_matrix (pd.DataFrame): A matrix showing barcode overlap counts between samples.
        pval_matrix (pd.DataFrame): A matrix showing corresponding p-values.
        significant_df (pd.DataFrame): Filtered significant sample pairs (P < threshold).
    """
    # Extract cell barcodes as sets
    barcode_sets = [set(adata.obs_names) for adata in adata_list]
    barcode_counts = [len(barcodes) for barcodes in barcode_sets]  # Sizes of each sample
    # Use sample names if provided, otherwise assign default names
    if sample_names is None:
        sample_names = [f"Sample_{i+1}" for i in range(len(adata_list))]
    # Create empty DataFrames to store overlap counts and p-values
    overlap_matrix = pd.DataFrame(index=sample_names, columns=sample_names, dtype=int)
    pval_matrix = pd.DataFrame(index=sample_names, columns=sample_names, dtype=float)
    significant_pairs = []
    # Compute pairwise overlaps and p-values
    for i, name1 in enumerate(sample_names):
        for j, name2 in enumerate(sample_names):
            if i < j:  # Avoid self-comparison and redundant calculations
                overlap_count = len(barcode_sets[i] & barcode_sets[j])
                overlap_matrix.loc[name1, name2] = overlap_count
                overlap_matrix.loc[name2, name1] = overlap_count  # Symmetric
                # Hypergeometric test
                N = total_barcodes
                K = barcode_counts[i]
                n = barcode_counts[j]
                k = overlap_count
                pval = stats.hypergeom.sf(k-1, N, K, n)  # Survival function (1 - CDF)
                pval_matrix.loc[name1, name2] = pval
                pval_matrix.loc[name2, name1] = pval  # Symmetric
                # Store significant pairs
                if pval < pval_threshold:
                    significant_pairs.append((name1, name2, overlap_count, pval))
    # Convert significant pairs to DataFrame
    significant_df = pd.DataFrame(significant_pairs, columns=["Sample_1", "Sample_2", "Overlap_Count", "P_Value"])
    return overlap_matrix, pval_matrix, significant_df

# Example execution (to be run with actual adata_list)
# overlap_df, pval_df, significant_df = compute_barcode_overlap_with_pval_and_filter(adata_list, sample_names)


import numpy as np
def compute_upper_triangle_pval_df(pval_df):
    """
    Extracts the upper triangular part of the p-value DataFrame (excluding the diagonal).

    Parameters:
        pval_df (pd.DataFrame): Symmetric DataFrame containing p-values.

    Returns:
        list: List of unique p-values from the upper triangle (excluding the diagonal).
    """
    upper_triangle_pvals = []
    sample_names = pval_df.index
    for i, name1 in enumerate(sample_names):
        for j, name2 in enumerate(sample_names):
            if i < j:  # Upper triangle only, excluding diagonal
                upper_triangle_pvals.append(pval_df.loc[name1, name2])
    return np.array(upper_triangle_pvals)



def plot_pval_histogram_upper_half_with_counts(pval_df, bins=50):
    """
    Plot a histogram of -log10 transformed p-values from the upper triangle of pval_df,
    displaying text only for bins with counts greater than 0.
    Parameters:
        pval_df (pd.DataFrame): DataFrame containing p-values.
        bins (int): Number of bins for the histogram.
    """
    # Extract upper triangular p-values
    upper_triangle_pvals = compute_upper_triangle_pval_df(pval_df)
    # Convert to -log10 scale, avoiding log(0) by setting a lower bound
    upper_triangle_pvals = np.clip(upper_triangle_pvals, 1e-300, 1)
    log_pval_values = -np.log10(upper_triangle_pvals)
    # Plot histogram
    plt.figure(figsize=(8, 6))
    counts, edges, patches = plt.hist(log_pval_values, bins=bins, edgecolor='black', alpha=0.75)
    # Annotate bars with counts if count > 0
    for count, edge, patch in zip(counts, edges[:-1], patches):
        if count > 0:
            plt.text(edge + (edges[1] - edges[0]) / 2, count, str(int(count)), 
                     ha='center', va='bottom', fontsize=10, color='black')
    plt.xlabel("-log10(P-value)")
    plt.ylabel("Frequency")
    plt.title("Histogram of -log10 Transformed P-values (Upper Triangle Only)")
    plt.grid(True, linestyle="--", linewidth=0.5)
    plt.show()

# Example execution (assuming pval_df is available)
# plot_pval_histogram_upper_half_with_counts(pval_df)



def plot_batch_fraction_boxplot(adata_merged, threshold = 5, figsize=(12, 3), dpi=150, save_path=None):
    batch_means = adata_merged.obs.groupby("Batch")["batch_purity_median"].median()
    # Compute the interquartile range (IQR) for batch means
    Q1 = batch_means.quantile(0.25)  # First quartile (25%)
    Q3 = batch_means.quantile(0.75)  # Third quartile (75%)
    IQR = Q3 - Q1  # Interquartile range
    #Define a soft threshold based on the IQR (e.g., considering values above Q3 + 1.5*IQR as high)
    upper_threshold = Q3 + threshold * IQR
    # Identify batches above the soft threshold
    high_fraction_batches = batch_means[batch_means > upper_threshold].sort_values(ascending=False)
    # Ensure correct annotation position by aligning `i` with the box positions
    plt.figure(figsize=figsize, dpi=dpi)
    ax = sns.boxplot(x=adata_merged.obs['Batch'], y=adata_merged.obs["batch_purity_median"], palette="Set2")
    # Find the maximum y-value for positioning annotations
    y_max = adata_merged.obs["batch_purity_median"].max()
    # Get the sorted unique batch names used in the x-axis of the boxplot
    batch_order = sorted(adata_merged.obs['Batch'].unique())
    # Add '*' annotation on top of significant batches and color x-labels in red
    xtick_labels = []
    for batch in batch_order:
        if batch in set(high_fraction_batches.index):
            x_color = 'red'  # Color significant batch labels in red
            xtick_labels.append(f"\033[31m{batch}\033[0m")  # ANSI escape code for red
            x_pos = batch_order.index(batch)  # Get the correct x-axis position
            ax.text(x_pos, y_max * 0.98, "*", ha='center', va='bottom', fontsize=14, fontweight='bold', color='red')
        else:
            x_color = 'black'
            xtick_labels.append(batch)
    # Set customized x-tick labels
    ax.set_xticklabels(batch_order, rotation=90, fontsize=10)
    # Customize plot appearance
    plt.xlabel("Batch", fontsize=12)
    plt.ylabel("Batch purity", fontsize=12)
    plt.title(f"Batch purity Across Batches (> Q3 + {threshold}*IQR)", fontsize=14)
    # Adjust x-tick colors
    for tick_label, batch in zip(ax.get_xticklabels(), batch_order):
        if batch in set(high_fraction_batches.index):
            tick_label.set_color('red')  # Set the text color of significant batch labels
    # Save the plot if a path is provided
    if save_path:
        plt.savefig(save_path, bbox_inches="tight", dpi=dpi)
        print(f"Plot saved to {save_path}")
    # Show the plot
    plt.show()



import numpy as np
import matplotlib.pyplot as plt
import math
from itertools import cycle

def plot_category_panels_same_space(
    adata, 
    category_column, 
    basis="umap", 
    max_cols=10, 
    figsize_base=(5, 5), 
    category_order=None,
    bg_color="#D9D9D9",        # 배경(전체 셀) 색
    fg_colors=None,            # 카테고리 강조 색 목록 (None이면 자동 순환)
    bg_size=2,                 # 배경 점 크기
    fg_size=4,                 # 강조 점 크기
    rasterized=True            # 래스터라이즈(큰 데이터에 유리)
):
    """
    같은 UMAP(또는 다른 embedding) 좌표계를 유지한 채 카테고리별 패널을 그립니다.
    각 패널: 전체 셀 회색 배경 + 해당 카테고리만 강조.
    """

    # 좌표 가져오기
    key = f"X_{basis}" if not basis.startswith("X_") else basis
    if key not in adata.obsm_keys():
        raise KeyError(f"'{key}' not found in adata.obsm. (e.g., use basis='umap' for 'X_umap')")

    coords = adata.obsm[key]
    x, y = coords[:, 0], coords[:, 1]

    # 축 범위 (모든 패널에서 동일하게 사용)
    xpad = (x.max() - x.min()) * 0.02
    ypad = (y.max() - y.min()) * 0.02
    xlim = (x.min() - xpad, x.max() + xpad)
    ylim = (y.min() - ypad, y.max() + ypad)

    # 카테고리 정리
    counts = adata.obs[category_column].value_counts()
    if category_order is not None:
        cats = [c for c in category_order if c in counts.index]
    else:
        cats = list(counts.index)
    n = len(cats)

    if n == 0:
        raise ValueError(f"No categories found in '{category_column}'.")

    # 그리드 계산
    cols = min(max_cols, n)
    rows = math.ceil(n / cols)

    # 색상 순환 준비
    if fg_colors is None:
        fg_colors = cycle(plt.rcParams["axes.prop_cycle"].by_key().get("color", ["C0"]))
    else:
        fg_colors = cycle(fg_colors)

    # Figure 생성
    fig, axs = plt.subplots(rows, cols, figsize=(figsize_base[0]*cols, figsize_base[1]*rows))
    if rows == 1 and cols == 1:
        axs = np.array([[axs]])
    elif rows == 1:
        axs = np.array([axs])
    elif cols == 1:
        axs = np.array([[ax] for ax in axs])

    # 카테고리별 패널 그리기
    for i, cat in enumerate(cats):
        r, c = divmod(i, cols)
        ax = axs[r, c]

        # 배경: 모든 셀
        ax.scatter(
            x, y, s=bg_size, c=bg_color, linewidths=0, rasterized=rasterized
        )

        # 포그라운드: 해당 카테고리만
        mask = (adata.obs[category_column] == cat).to_numpy()
        col = next(fg_colors)
        ax.scatter(
            x[mask], y[mask], s=fg_size, c=col, linewidths=0, rasterized=rasterized
        )

        # 동일 축 범위/비율 고정
        ax.set_xlim(*xlim)
        ax.set_ylim(*ylim)
        ax.set_aspect("equal", adjustable="box")
        ax.set_xticks([]); ax.set_yticks([])
        ax.set_frame_on(False)
        ax.set_title(f"{cat} (n={int(mask.sum())})", fontsize=11)

    # 남는 패널 비활성화
    for j in range(n, rows*cols):
        r, c = divmod(j, cols)
        axs[r, c].axis("off")

    plt.tight_layout()
    plt.show()
    return fig







import matplotlib
import seaborn as sns
import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
import numpy as np
import pandas as pd
import scanpy as sc
from adjustText import adjust_text  # Import adjustText for label adjustment

sns.set_style('ticks')
matplotlib.rcParams['figure.figsize'] = [6, 6]
matplotlib.rcParams['figure.dpi'] = 200

def plot_umap_with_centroids(adata, color, annotate=True, figsize=matplotlib.rcParams['figure.figsize'], dpi=matplotlib.rcParams['figure.dpi']):
    """
    Plots a UMAP projection with optional centroids annotated for the specified category.
    Ensures proper categorical coloring and prevents text overlap.
    
    Parameters:
    - adata: AnnData object containing UMAP coordinates and category annotations.
    - color: str, name of the column in `adata.obs` for annotations (e.g., 'Celltype').
    - annotate: bool, whether to annotate centroids on the plot (default: True).
    - figsize: tuple, size of the figure (default: (4, 4)).
    - dpi: int, resolution of the figure in dots per inch (default: 200).
    """
    # Convert category to categorical type if not already
    if not isinstance(adata.obs[color].dtype, pd.CategoricalDtype):
        adata.obs[color] = adata.obs[color].astype("category")
    # Generate distinct colors for all unique categories
    num_categories = adata.obs[color].nunique()
    palette = sns.color_palette("hsv", num_categories)  # Use HSV to ensure color distinction
    adata.uns[f"{color}_colors"] = [mcolors.rgb2hex(color) for color in palette]
    # Plot UMAP without legend
    sc.pl.umap(
        adata,
        color=[color],
        frameon=True,
        ncols=1,
        show=False,
        legend_loc=None,
        size=1
    )
    if annotate:
        # Calculate centroids
        umap_coordinates = adata.obsm['X_umap']
        categories = adata.obs[color]
        centroids = {cat: umap_coordinates[categories == cat].mean(axis=0) for cat in np.unique(categories)}
        # Annotate centroids on the UMAP plot with dynamic text adjustment
        ax = plt.gca()
        texts = [ax.text(centroid[0], centroid[1], cat, fontsize=10, ha='center', va='center', fontweight="bold")
                 for cat, centroid in centroids.items()]
        adjust_text(texts, ax=ax, arrowprops=dict(arrowstyle="-", lw=0.5))  # Adjust text to prevent overlap
    plt.show()







def plot_category_panels(adata, category_column, basis="umap", max_cols=10, figsize_base=(5, 5)):
    """
    Plots a figure panel across each category in the specified column, 
    showing embeddings colored by the category.
    
    Parameters:
    - adata: AnnData object
    - category_column: str, column in adata.obs to group by
    - basis: str, embedding basis for plotting (default: "X_mde")
    - max_cols: int, maximum number of columns in the figure grid (default: 4)
    - figsize_base: tuple, base figure size per subplot (default: (5,5))
    
    Returns:
    - A matplotlib figure with panels for each category
    """
    # Get unique categories and count the number of items per category
    category_counts = adata.obs[category_column].value_counts()
    unique_categories = category_counts.index
    num_categories = len(unique_categories)
    # Determine rows and columns for subplot grid
    cols = min(max_cols, num_categories)  
    rows = (num_categories // max_cols) + (num_categories % max_cols > 0)  # Adjust rows dynamically
    # Create figure with dynamic size
    fig, axs = plt.subplots(rows, cols, figsize=(figsize_base[0] * cols, figsize_base[1] * rows))
    # Ensure axs is always a 2D array
    if rows == 1:
        axs = axs.reshape(1, -1)
    elif cols == 1:
        axs = axs.reshape(-1, 1)
    # Plot each category
    for idx, category in enumerate(unique_categories):
        ax = axs[idx // max_cols, idx % max_cols]
        count = category_counts[category]
        sc.pl.embedding(
            adata[adata.obs[category_column] == category, :], 
            basis=basis, 
            color=category_column,
            ax=ax, 
            show=False,
            title=f"{category} (n={count})",
            frameon=False
        )
    # Turn off any excess subplots
    for idx in range(num_categories, rows * cols):
        axs[idx // max_cols, idx % max_cols].axis('off')
    plt.tight_layout()
    plt.show()





### dup
#FY24/Br6312_DLPFC_A1 by FY24/Br1390_Hippo_G10
#/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY24/Br1390_Hippo_G10
#adata = sc.read_10x_h5("/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY24/Br6312_DLPFC_A1/cellbender_output_filtered.h5")
#bdata = sc.read_10x_h5("/mnt/pv_compute/dongsan/datasets/snRNA-seq/AANRI_claust_brain_placenta/AANRI_36-38_n12_NS2K_05172024/Br1390_Hippo_G10/cellbender_output_filtered.h5")
#adata_unique = adata[~adata.obs.index.isin(bdata.obs.index)].copy()
#adata_unique.write_h5ad("/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY24/Br6312_DLPFC_A1/cellbender_output_filtered_dup_removed.h5ad")

#adata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY23/Br5338_hippo_A9/cellbender_output_filtered.h5')
#bdata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY23/Br5562_caudate_A8/cellbender_output_filtered.h5')
#sc.pp.filter_cells(adata,300)
#sc.pp.filter_cells(bdata,300)
#inter = set(adata.obs.index).intersection(bdata.obs.index)
#aa=adata.obs[adata.obs.index.isin(inter)]
#bb=bdata.obs[bdata.obs.index.isin(inter)]
#aa["group"] = "A"
#bb["group"] = "B"
#df = pd.concat([aa, bb])
#sns.boxenplot(x="group", y="n_counts", data=df)
#plt.show()
#adata_unique = adata[~adata.obs.index.isin(bdata.obs.index)].copy()
#adata_unique.write_h5ad("/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY23/Br5338_hippo_A9/cellbender_output_filtered_dup_removed.h5ad")


#adata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY23/Br0863_hippo_C6/cellbender_output_filtered.h5')
#bdata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY23/Br0989_hippo_B12/cellbender_output_filtered.h5')
#sc.pp.filter_cells(adata,300)
#sc.pp.filter_cells(bdata,300)
#inter = set(adata.obs.index).intersection(bdata.obs.index)
#aa=adata.obs[adata.obs.index.isin(inter)]
#bb=bdata.obs[bdata.obs.index.isin(inter)]
#aa['group'] = 'A'
#bb['group'] = 'B'
#df = pd.concat([aa, bb])
#sns.boxenplot(x="group", y="n_counts", data=df)
#plt.show()
#print(len(inter))
#adata_unique = adata[~adata.obs.index.isin(bdata.obs.index)].copy()
#adata_unique.write_h5ad("/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY23/Br0863_hippo_C6/cellbender_output_filtered_dup_removed.h5ad")


#adata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY24/Br1135_Hippo_F4/cellbender_output_filtered.h5')
#bdata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY24/Br1227_DLPFC_E8/cellbender_output_filtered.h5')
#sc.pp.filter_cells(adata,300)
#sc.pp.filter_cells(bdata,300)
#inter = set(adata.obs.index).intersection(bdata.obs.index)
#aa=adata.obs[adata.obs.index.isin(inter)]
#bb=bdata.obs[bdata.obs.index.isin(inter)]
#aa['group'] = 'A'
#bb['group'] = 'B'
#df = pd.concat([aa, bb])
#sns.boxenplot(x="group", y="n_counts", data=df)
#plt.show()
#print(len(inter))
#adata_unique = adata[~adata.obs.index.isin(bdata.obs.index)].copy()
#adata_unique.write_h5ad("/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY24/Br1135_Hippo_F4/cellbender_output_filtered_dup_removed.h5ad")

#adata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY24/Br1442_Hippo_F2/cellbender_output_filtered.h5')
#bdata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY24/Br1227_DLPFC_E8/cellbender_output_filtered.h5')
#sc.pp.filter_cells(adata,300)
#sc.pp.filter_cells(bdata,300)
#inter = set(adata.obs.index).intersection(bdata.obs.index)
#aa=adata.obs[adata.obs.index.isin(inter)]
#bb=bdata.obs[bdata.obs.index.isin(inter)]
#aa['group'] = 'A'
#bb['group'] = 'B'
#df = pd.concat([aa, bb])
#sns.boxenplot(x="group", y="n_counts", data=df)
#plt.show()
#print(len(inter))
#adata_unique = adata[~adata.obs.index.isin(bdata.obs.index)].copy()
#adata_unique.write_h5ad("/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY24/Br1442_Hippo_F2/cellbender_output_filtered_dup_removed.h5ad")

#adata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY24/Br1722_Hippo_D12/cellbender_output_filtered.h5')
#bdata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY24/Br6120_Hippo_E1/cellbender_output_filtered.h5')
#sc.pp.filter_cells(adata,300)
#sc.pp.filter_cells(bdata,300)
#inter = set(adata.obs.index).intersection(bdata.obs.index)
#aa=adata.obs[adata.obs.index.isin(inter)]
#bb=bdata.obs[bdata.obs.index.isin(inter)]
#aa['group'] = 'A'
#bb['group'] = 'B'
#df = pd.concat([aa, bb])
#sns.boxenplot(x="group", y="n_counts", data=df)
#plt.show()
#print(len(inter))
#adata_unique = adata[~adata.obs.index.isin(bdata.obs.index)].copy()
#adata_unique.write_h5ad("/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY24/Br1722_Hippo_D12/cellbender_output_filtered_dup_removed.h5ad")


#adata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY23/Br5370_caudate_E11/cellbender_output_filtered.h5')
#bdata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY23/Br1539_hippo_A7/cellbender_output_filtered.h5')
#sc.pp.filter_cells(adata,300)
#sc.pp.filter_cells(bdata,300)
#inter = set(adata.obs.index).intersection(bdata.obs.index)
#aa=adata.obs[adata.obs.index.isin(inter)]
#bb=bdata.obs[bdata.obs.index.isin(inter)]
#aa['group'] = 'A'
#bb['group'] = 'B'
#df = pd.concat([aa, bb])
#sns.boxenplot(x="group", y="n_counts", data=df)
#plt.show()
#print(len(inter))
#adata_unique = adata[~adata.obs.index.isin(bdata.obs.index)].copy()
#adata_unique.write_h5ad("/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY23/Br5370_caudate_E11/cellbender_output_filtered_dup_removed.h5ad")

#adata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY25/Br6129_DLPFC_G1/cellbender_output_filtered.h5')
#bdata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY25/Br1571_Hippo_E12/cellbender_output_filtered.h5')
#sc.pp.filter_cells(adata,300)
#sc.pp.filter_cells(bdata,300)
#inter = set(adata.obs.index).intersection(bdata.obs.index)
#aa=adata.obs[adata.obs.index.isin(inter)]
#bb=bdata.obs[bdata.obs.index.isin(inter)]
#aa['group'] = 'A'
#bb['group'] = 'B'
#df = pd.concat([aa, bb])
#sns.boxenplot(x="group", y="n_counts", data=df)
#plt.show()
#print(len(inter))
#adata_unique = adata[~adata.obs.index.isin(bdata.obs.index)].copy()
#adata_unique.write_h5ad("/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY25/Br6129_DLPFC_G1/cellbender_output_filtered_dup_removed.h5ad")

#adata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY25/Br1227_DLPFC_E8/cellbender_output_filtered.h5')
#bdata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY25/Br1611_DLPFC_C4/cellbender_output_filtered.h5')
#sc.pp.filter_cells(adata,300)
#sc.pp.filter_cells(bdata,300)
#inter = set(adata.obs.index).intersection(bdata.obs.index)
#aa=adata.obs[adata.obs.index.isin(inter)]
#bb=bdata.obs[bdata.obs.index.isin(inter)]
#aa['group'] = 'A'
#bb['group'] = 'B'
#df = pd.concat([aa, bb])
#sns.boxenplot(x="group", y="n_counts", data=df)
#plt.show()
#print(len(inter))
#adata_unique = adata[~adata.obs.index.isin(bdata.obs.index)].copy()
#adata_unique.write_h5ad("/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY25/Br1185_DLPFC_C11/cellbender_output_filtered_dup_removed.h5ad")

#adata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY25/Br1473_Hippo_G4/cellbender_output_filtered.h5')
#bdata = sc.read_10x_h5('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY25/Br1571_Hippo_E12/cellbender_output_filtered.h5')
#sc.pp.filter_cells(adata,300)
#sc.pp.filter_cells(bdata,300)
#inter = set(adata.obs.index).intersection(bdata.obs.index)
#aa=adata.obs[adata.obs.index.isin(inter)]
#bb=bdata.obs[bdata.obs.index.isin(inter)]
#aa['group'] = 'A'
#bb['group'] = 'B'
#df = pd.concat([aa, bb])
#sns.boxenplot(x="group", y="n_counts", data=df)
#plt.show()
#print(len(inter))
#adata_unique = adata[~adata.obs.index.isin(bdata.obs.index)].copy()
#adata_unique.write_h5ad("/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY25/Br1473_Hippo_G4/cellbender_output_filtered_dup_removed.h5ad")





#	nova_id	            region	BrNum
#29	Br1275_hippo_C2	    caudate	Br1275  annotated hippo but caudate... (next nova ID unmatch...)
#45	Br1435_caudate_C1	hippo	Br1435  annotated caudate but hippo...
#69	Br5338_Caudate_E01	caudate	Br5338
#93	Br8181_Caudate_D12	caudate	Br8181
#95	Br8181_Hippo_D11	hippo	Br8181

#"Br1410_FY23 needs to be rerun." excluded due to low umap quality... (too many cells with unstructured clusters)






import glob

Region = 'hippo'

if Region == 'caudate':
    path="/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY23/"
    subdirs = [d for d in os.listdir(path) if os.path.isdir(os.path.join(path, d))]
    matching_files_FY23 = [d for d in subdirs if "audate" in d]
    matching_files_FY23 = [os.path.join(path, d) for d in matching_files_FY23]
    path="/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY24/"
    subdirs = [d for d in os.listdir(path) if os.path.isdir(os.path.join(path, d))]
    matching_files_FY24 = [d for d in subdirs if "audate" in d]
    matching_files_FY24 = [os.path.join(path, d) for d in matching_files_FY24]
    path="/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY25/"
    subdirs = [d for d in os.listdir(path) if os.path.isdir(os.path.join(path, d))]
    matching_files_FY25 = [d for d in subdirs if "audate" in d]
    matching_files_FY25 = [os.path.join(path, d) for d in matching_files_FY25]
    selected_files = []
    for folder in matching_files_FY23:
        h5ad_path = os.path.join(folder, "cellbender_output_filtered_dup_removed.h5ad")
        h5_path   = os.path.join(folder, "cellbender_output_filtered.h5")
        if os.path.exists(h5ad_path):
            selected_files.append(h5ad_path)
        elif os.path.exists(h5_path):
            selected_files.append(h5_path)
    matching_files_FY23 = selected_files
    selected_files = []
    for folder in matching_files_FY24:
        h5ad_path = os.path.join(folder, "cellbender_output_filtered_dup_removed.h5ad")
        h5_path   = os.path.join(folder, "cellbender_output_filtered.h5")
        if os.path.exists(h5ad_path):
            selected_files.append(h5ad_path)
        elif os.path.exists(h5_path):
            selected_files.append(h5_path)
    matching_files_FY24 = selected_files
    selected_files = []
    for folder in matching_files_FY25:
        h5ad_path = os.path.join(folder, "cellbender_output_filtered_dup_removed.h5ad")
        h5_path   = os.path.join(folder, "cellbender_output_filtered.h5")
        if os.path.exists(h5ad_path):
            selected_files.append(h5ad_path)
        elif os.path.exists(h5_path):
            selected_files.append(h5_path)
    matching_files_FY25 = selected_files   
    matching_files = matching_files_FY23 + matching_files_FY24 + matching_files_FY25
    matching_files.append('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY23/Br1275_hippo_C2/cellbender_output_filtered.h5') 
    matching_files = [item for item in matching_files if not any(exclude in item for exclude in ["Br1435_caudate_C1", "Br0846_caudate_E12","Br1285_Caudate_YKL3_A2","Br1137_Caudate_C10","Br1918_Caudate_C1","Br0846_Caudate_D4","Br1918_Caudate_E2","Br2623_Caudate_D7","Br0989_caudate_G2"])] #Br1435_caudate_C1 : hippo

if Region == 'DLPFC':
    path="/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY23/"
    subdirs = [d for d in os.listdir(path) if os.path.isdir(os.path.join(path, d))]
    matching_files_FY23 = [d for d in subdirs if "DLPFC" in d]
    matching_files_FY23 = [os.path.join(path, d) for d in matching_files_FY23]
    path="/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY24/"
    subdirs = [d for d in os.listdir(path) if os.path.isdir(os.path.join(path, d))]
    matching_files_FY24 = [d for d in subdirs if "DLPFC" in d]
    matching_files_FY24 = [os.path.join(path, d) for d in matching_files_FY24]
    path="/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY25/"
    subdirs = [d for d in os.listdir(path) if os.path.isdir(os.path.join(path, d))]
    matching_files_FY25 = [d for d in subdirs if "DLPFC" in d]
    matching_files_FY25 = [os.path.join(path, d) for d in matching_files_FY25]
    selected_files = []
    for folder in matching_files_FY23:
        h5ad_path = os.path.join(folder, "cellbender_output_filtered_dup_removed.h5ad")
        h5_path   = os.path.join(folder, "cellbender_output_filtered.h5")
        if os.path.exists(h5ad_path):
            selected_files.append(h5ad_path)
        elif os.path.exists(h5_path):
            selected_files.append(h5_path)
    matching_files_FY23 = selected_files
    matching_files_FY23 = [item for item in matching_files_FY23 if "Br0982" not in item]
    selected_files = []
    for folder in matching_files_FY24:
        h5ad_path = os.path.join(folder, "cellbender_output_filtered_dup_removed.h5ad")
        h5_path   = os.path.join(folder, "cellbender_output_filtered.h5")
        if os.path.exists(h5ad_path):
            selected_files.append(h5ad_path)
        elif os.path.exists(h5_path):
            selected_files.append(h5_path)
    matching_files_FY24 = selected_files
    matching_files_FY24 = [item for item in matching_files_FY24 if "Br2388" not in item] 
    selected_files = []
    for folder in matching_files_FY25:
        h5ad_path = os.path.join(folder, "cellbender_output_filtered_dup_removed.h5ad")
        h5_path   = os.path.join(folder, "cellbender_output_filtered.h5")
        if os.path.exists(h5ad_path):
            selected_files.append(h5ad_path)
        elif os.path.exists(h5_path):
            selected_files.append(h5_path)
    matching_files_FY25 = selected_files   
    matching_files = matching_files_FY23 + matching_files_FY24 + matching_files_FY25
    matching_files = [item for item in matching_files if not any(exclude in item for exclude in ["Br1410_DLPFC_D8", "Br1442_DLPFC_E12", "Br1342_DLPFC_F2","Br1342_DLPFC_D1","Br1193_DLPFC_B2","Br1137_DLPFC_D10","Br1137_DLPFC_C8"])] 

if Region == 'hippo':
    path="/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY23/"
    subdirs = [d for d in os.listdir(path) if os.path.isdir(os.path.join(path, d))]
    matching_files_FY23 = [d for d in subdirs if "ippo" in d]
    matching_files_FY23 = [os.path.join(path, d) for d in matching_files_FY23]
    path="/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY24/"
    subdirs = [d for d in os.listdir(path) if os.path.isdir(os.path.join(path, d))]
    matching_files_FY24 = [d for d in subdirs if "ippo" in d]
    matching_files_FY24 = [os.path.join(path, d) for d in matching_files_FY24]
    path="/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY25/"
    subdirs = [d for d in os.listdir(path) if os.path.isdir(os.path.join(path, d))]
    matching_files_FY25 = [d for d in subdirs if "ippo" in d]
    matching_files_FY25 = [os.path.join(path, d) for d in matching_files_FY25]
    selected_files = []
    for folder in matching_files_FY23:
        h5ad_path = os.path.join(folder, "cellbender_output_filtered_dup_removed.h5ad")
        h5_path   = os.path.join(folder, "cellbender_output_filtered.h5")
        if os.path.exists(h5ad_path):
            selected_files.append(h5ad_path)
        elif os.path.exists(h5_path):
            selected_files.append(h5_path)
    matching_files_FY23 = selected_files
    matching_files_FY23.append('/mnt/pv_compute/dongsan/datasets/AANRI/NOVASEQ_FY23/Br1435_caudate_C1/cellbender_output_filtered.h5')  ## actually Hippo!
    selected_files = []
    for folder in matching_files_FY24:
        h5ad_path = os.path.join(folder, "cellbender_output_filtered_dup_removed.h5ad")
        h5_path   = os.path.join(folder, "cellbender_output_filtered.h5")
        if os.path.exists(h5ad_path):
            selected_files.append(h5ad_path)
        elif os.path.exists(h5_path):
            selected_files.append(h5_path)
    matching_files_FY24 = selected_files
    matching_files_FY24 = [item for item in matching_files_FY24 if "Br2388" not in item] 
    selected_files = []
    for folder in matching_files_FY25:
        h5ad_path = os.path.join(folder, "cellbender_output_filtered_dup_removed.h5ad")
        h5_path   = os.path.join(folder, "cellbender_output_filtered.h5")
        if os.path.exists(h5ad_path):
            selected_files.append(h5ad_path)
        elif os.path.exists(h5_path):
            selected_files.append(h5_path)
    matching_files_FY25 = selected_files   
    matching_files = matching_files_FY23 + matching_files_FY24 + matching_files_FY25
    matching_files = [item for item in matching_files if not any(exclude in item for exclude in ["Br1275_hippo_C2", "Br1324_Hippo_A3","Br5253_Hippo_D8","Br2623_Hippo_B4","Br1137_Hippo_G3","Br0991_Hippo_A1","Br0991_Hippo_YKL3_A3","Br1185_Hippo_H12","Br1297_Hippo_C9","Br1297_Hippo_D12","Br1324_Hippo_D2","Br1473_Hippo_D6","Br5253_Hippo_D3","Br2623_Hippo_E8","Br1522_Hippo_E11"])] 

if Region=='LIBD':
    sampleInfo = pd.read_excel('/mnt/pv_compute/dongsan/datasets/AANRI/LIBD_DLPFC/LIBD_DLPFC_CAUC_info.xlsx')
    filtered = sampleInfo[sampleInfo.PrimaryDx=='Control'].copy()
    filtered['SampleGroup'] = filtered['SampleID'].apply(
        lambda x: 'Punzi' if 'Punzi' in x else 'Else'
    )
    filtered['SampleGroup'] = filtered['SampleGroup'].astype('category')
    path = "/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered/analysis_per_sample.R"
    matching_files=[]
    for i in range(0,len(filtered.SampleID.tolist())):
        #print(f"/mnt/pv_compute/dongsan/datasets/AANRI/LIBD_DLPFC/cellranger8/{filtered.iloc[i].SampleID}/cellbender_output_filtered.h5")
        matching_files.append(f"/mnt/pv_compute/dongsan/datasets/AANRI/LIBD_DLPFC/cellranger8/{filtered.iloc[i].SampleID}/cellbender_output_filtered.h5")



data = []
for path in matching_files:
    parts = path.split('/')  # Split by "/"
    fy = parts[-3]  # FY24 (second last folder)
    sample_name =  parts[-2].replace(".h5ad", "")  # Extract sample name
    data.append((fy, sample_name))


if Region=='LIBD':
    matching_files_filtered=[]
    path = "/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/analysis_per_sample.R"
    sampleInfo = pd.read_excel('/mnt/pv_compute/dongsan/datasets/AANRI/LIBD_DLPFC/LIBD_DLPFC_CAUC_info.xlsx')
    global_ancestry = pd.read_excel('/mnt/pv_compute/dongsan/datasets/AANRI/global_ancestry.xlsx')
    print("Matching files:")
    for i in range(0,len(matching_files)):
        print(matching_files[i])
        ext = os.path.splitext(matching_files[i])[1].lower()
        if ext==".h5":
            adata = sc.read_10x_h5(matching_files[i])
        elif ext==".h5ad":
            adata = sc.read_h5ad(matching_files[i])
        filename = data[i][1]
        samp = sampleInfo.loc[sampleInfo.SampleID==filename]
        adata.obs['BrNum'] = samp['BrNum'].astype(str).iloc[0]
        adata.obs['AgeDeath'] = samp['AgeDeath'].iloc[0]
        adata.obs['Sex'] = samp['Sex'].astype(str).iloc[0]
        adata.obs['PMI'] = samp['PMI'].astype(str).iloc[0]
        adata.obs['Batch'] = filename  
        #adata.obs['AgeGroup'] = sampleInfo[sampleInfo.BrNum == adata.obs.BrNum.unique()[0]].AgeGroup.values[0]
        #adata.obs['YRI'] = global_ancestry[global_ancestry.BrNum == adata.obs.BrNum.unique()[0]].YRI.values[0]
        #adata.obs['CEU'] = global_ancestry[global_ancestry.BrNum == adata.obs.BrNum.unique()[0]].CEU.values[0]
        #adata.obs['CHB'] = global_ancestry[global_ancestry.BrNum == adata.obs.BrNum.unique()[0]].CHB.values[0]
        adata.obs['Region'] =  'DLPFC'
        adata.var_names_make_unique()
        sc.pp.filter_genes(adata, min_cells=3)
        sc.pp.filter_cells(adata, min_genes=400)
        # mitochondrial genes
        adata.var["mt"] = adata.var_names.str.startswith("MT-")
        # ribosomal genes
        adata.var["ribo"] = adata.var_names.str.startswith(("RPS", "RPL"))
        ## percent_top N: the occupancy of the most abundant N genes (a measure for gene expression diversity per cell)
        sc.pp.calculate_qc_metrics(adata, qc_vars=["mt", "ribo"], inplace=True, percent_top=[20,50,100], log1p=True)
        #adata.obs["mt_outlier"] = is_outlier(adata, "pct_counts_mt", 5) | (
        #adata.obs["pct_counts_mt"] > 10)
        print(f"Total number of cells before filtration: {adata.n_obs}")
        adata.obs["mt_outlier"] = adata.obs["pct_counts_mt"] > 5
        adata.obs["ribo_outlier"] = adata.obs["pct_counts_ribo"] > 10
        adata = adata[~adata.obs["mt_outlier"] & ~adata.obs["ribo_outlier"]].copy()
        sc.pp.scrublet(adata)
        adata = adata[(~adata.obs.predicted_doublet)].copy()   
        adata.obs.index = adata.obs.index.astype(str) + "_" + adata.obs["Batch"].astype(str)
        adata_raw = adata.copy()
        adata_raw.obs['predicted_doublet'] = adata.obs['predicted_doublet']
        adata_raw.obs['doublet_score'] = adata.obs['doublet_score']
        #model_seacell, adata = run_SEACells(adata)
        #adata_raw.obs['SEACell'] = adata.obs.Batch.astype(str) + '|' + adata.obs.SEACell.astype(str)
        adata_raw.obsp = adata.obsp.copy()
        adata_raw.obsm = adata.obsm.copy()
        adata_raw.uns = adata.uns.copy()
        adata_raw.write(f"/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/{filename}.h5ad")
        result=subprocess.run(["Rscript", path] + [Region] + [data[i][1]], capture_output=True, text=True)
        meta=pd.read_csv("/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/meta.csv",index_col=0)
        adata_raw.obs['Celltype_MTG'] = meta.Celltype_MTG
        #adata_raw.obs['Celltype'] = adata_raw.obs['Celltype_MTG'].copy()
        #SEACell_ad = SEACells.core.summarize_by_SEACell(adata_raw, SEACells_label='SEACell', summarize_layer='X')
        #SEACell_purity = SEACells.evaluate.compute_celltype_purity(adata_raw, 'Celltype')
        #SEACell_ad.obs['Celltype'] = SEACell_purity['Celltype'] 
        #SEACell_ad.obs['Celltype_purity'] = SEACell_purity['Celltype_purity']
        #sc.pp.normalize_total(SEACell_ad)
        #sc.pp.log1p(SEACell_ad)
        #sc.pp.highly_variable_genes(SEACell_ad, n_top_genes=3000)
        #sc.tl.pca(SEACell_ad, n_comps=50, use_highly_variable=True)
        #sc.pp.neighbors(SEACell_ad)
        #sc.tl.umap(SEACell_ad)
        #SEACell_ad.write('/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/'+ data[i][1] +'_SEACell_summarized.h5ad')
        adata_raw.write(f"/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/{filename}.h5ad")
        matching_files_filtered.append(f"/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/{filename}.h5ad")
        print(f"Total number of cells after filtration: {adata_raw.n_obs}")
else:
    path = "/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/analysis_per_sample.R"
    sampleInfo = pd.read_excel('/mnt/pv_compute/dongsan/datasets/AANRI/sample_inf.xlsx')
    global_ancestry = pd.read_excel('/mnt/pv_compute/dongsan/datasets/AANRI/global_ancestry.xlsx')
    matching_files_filtered=[]
    print("Matching files:")
    for i in range(0,len(matching_files)):
        print(matching_files[i])
        ext = os.path.splitext(matching_files[i])[1].lower()
        if ext==".h5":
            adata = sc.read_10x_h5(matching_files[i])
        elif ext==".h5ad":
            adata = sc.read_h5ad(matching_files[i])
        filename = data[i][1]
        if filename == "Br0948_Caudate_B9":
            filename = "Br0991_Caudate_B9"
        match = re.search(r'(Br\d+)', filename)
        adata.obs['Batch'] = filename
        adata.obs['BrNum'] = match.group(1)
        adata.obs['AgeDeath'] = sampleInfo[sampleInfo.BrNum == adata.obs.BrNum.unique()[0]].AgeDeath.values[0]
        adata.obs['Sex'] = sampleInfo[sampleInfo.BrNum == adata.obs.BrNum.unique()[0]].Sex.values[0]
        adata.obs['PMI'] = sampleInfo[sampleInfo.BrNum == adata.obs.BrNum.unique()[0]].PMI.values[0]
        adata.obs['AgeGroup'] = sampleInfo[sampleInfo.BrNum == adata.obs.BrNum.unique()[0]].AgeGroup.values[0]
        adata.obs['YRI'] = global_ancestry[global_ancestry.BrNum == adata.obs.BrNum.unique()[0]].YRI.values[0]
        adata.obs['CEU'] = global_ancestry[global_ancestry.BrNum == adata.obs.BrNum.unique()[0]].CEU.values[0]
        adata.obs['CHB'] = global_ancestry[global_ancestry.BrNum == adata.obs.BrNum.unique()[0]].CHB.values[0]
        adata.obs['Region'] =  Region
        adata.var_names_make_unique()
        sc.pp.filter_genes(adata, min_cells=3)
        sc.pp.filter_cells(adata, min_genes=400)
        # mitochondrial genes
        adata.var["mt"] = adata.var_names.str.startswith("MT-")
        # ribosomal genes
        adata.var["ribo"] = adata.var_names.str.startswith(("RPS", "RPL"))
        ## percent_top N: the occupancy of the most abundant N genes (a measure for gene expression diversity per cell)
        sc.pp.calculate_qc_metrics(adata, qc_vars=["mt", "ribo"], inplace=True, percent_top=[20,50,100], log1p=True)
        #adata.obs["mt_outlier"] = is_outlier(adata, "pct_counts_mt", 5) | (
        #adata.obs["pct_counts_mt"] > 10)
        print(f"Total number of cells before filtration: {adata.n_obs}")
        adata.obs["mt_outlier"] = adata.obs["pct_counts_mt"] > 5
        adata.obs["ribo_outlier"] = adata.obs["pct_counts_ribo"] > 10
        adata = adata[~adata.obs["mt_outlier"] & ~adata.obs["ribo_outlier"]].copy()
        sc.pp.scrublet(adata)
        adata = adata[(~adata.obs.predicted_doublet)].copy()   
        adata.obs.index = adata.obs.index.astype(str) + "_" + adata.obs["Batch"].astype(str)
        adata_raw = adata.copy()
        adata_raw.obs['predicted_doublet'] = adata.obs['predicted_doublet']
        adata_raw.obs['doublet_score'] = adata.obs['doublet_score']
        #model_seacell, adata = run_SEACells(adata)
        #adata_raw.obs['SEACell'] = adata.obs.Batch.astype(str) + '|' + adata.obs.SEACell.astype(str)
        adata_raw.obsp = adata.obsp.copy()
        adata_raw.obsm = adata.obsm.copy()
        adata_raw.uns = adata.uns.copy()
        adata_raw.write(f"/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/{filename}.h5ad")
        result=subprocess.run(["Rscript", path] + [Region] + [data[i][1]], capture_output=True, text=True)
        meta=pd.read_csv("/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/meta.csv",index_col=0)
        if Region=='caudate':
            adata_raw.obs['Celltype_caudate_NCOMM'] = meta.Celltype_caudate_NCOMM
            adata_raw.obs['Celltype_caudate_SCIENCE'] = meta.Celltype_caudate_SCIENCE
        elif Region=='DLPFC':
            adata_raw.obs['Celltype_MTG'] = meta.Celltype_MTG
        elif Region=='hippo':
            adata_raw.obs['Celltype_hippo_LIBD'] = meta.Celltype_hippo_LIBD
            adata_raw.obs['Celltype_hippo_SCIENCE'] = meta.Celltype_hippo_SCIENCE
            #adata_raw.obs['Celltype'] = adata_raw.obs['Celltype_MTG'].copy()
        #SEACell_ad = SEACells.core.summarize_by_SEACell(adata_raw, SEACells_label='SEACell', summarize_layer='X')
        #SEACell_purity = SEACells.evaluate.compute_celltype_purity(adata_raw, 'Celltype')
        #SEACell_ad.obs['Celltype'] = SEACell_purity['Celltype'] 
        #SEACell_ad.obs['Celltype_purity'] = SEACell_purity['Celltype_purity']
        #sc.pp.normalize_total(SEACell_ad)
        #sc.pp.log1p(SEACell_ad)
        #sc.pp.highly_variable_genes(SEACell_ad, n_top_genes=3000)
        #sc.tl.pca(SEACell_ad, n_comps=50, use_highly_variable=True)
        #sc.pp.neighbors(SEACell_ad)
        #sc.tl.umap(SEACell_ad)
        #SEACell_ad.write('/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/'+ data[i][1] +'_SEACell_summarized.h5ad')
        adata_raw.write(f"/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/{filename}.h5ad")
        matching_files_filtered.append(f"/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/{filename}.h5ad")
        print(f"Total number of cells after filtration: {adata_raw.n_obs}")


with open('/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/'+Region+'.txt', "w") as f:
    f.write("\n".join(matching_files_filtered))


import os
os.environ["SCIPY_ARRAY_API"] = "1"
from sklearn.metrics import classification_report
from scarches.models.scpoli import scPoli
import scipy.sparse as sp
from anndata import AnnData
from scipy.sparse import issparse



def run_scpoli(scpoli_model, source_adata, target_adata, celltype):
    early_stopping_kwargs = {
    "early_stopping_metric": "val_prototype_loss",
    "mode": "min",
    "threshold": 0,
    "patience": 20,
    "reduce_lr": True,
    "lr_patience": 13,
    "lr_factor": 0.1}
    X = target_adata.X.toarray() if issparse(target_adata.X) else target_adata.X
    df = pd.DataFrame(X, index=target_adata.obs_names, columns=target_adata.var_names)
    df = df.reindex(columns=source_adata.var_names, fill_value=0)        
    new_adata = AnnData(X=df.values.astype(np.float32), obs=target_adata.obs.copy(), var=source_adata.var.copy())
    adata = new_adata.copy()
    adata.obs['celltype_annotated'] = adata.obs[celltype].copy()
    adata.obs['conditions_combined'] = adata.obs.Batch.copy()
    adata.obs['Sample'] = adata.obs.Batch.copy()
    empty_idx = np.array([], dtype=int)
    ## initial celltype label transfer
    scpoli_query = scPoli.load_query_data(adata=adata,reference_model=scpoli_model,labeled_indices=empty_idx)
    scpoli_query.train(n_epochs=200,pretraining_epochs=160,eta=5)
    results_dict = scpoli_query.classify(adata, scale_uncertainties=True)
    adata.obs['Celltype_scpoli'] = results_dict['celltype_annotated']["preds"]
    adata.obs['Celltype_scpoli_uncert'] = results_dict['celltype_annotated']["uncert"]
    a_raw = adata.obs['Celltype_scpoli'].astype("string")
    b_raw = adata.obs[celltype].astype("string")
    # 2) 일치 마스크(둘 다 결측 아니고, 정규화 후 동일)
    agree = a_raw.notna() & b_raw.notna() & (a_raw == b_raw)
    # 3) 합의 라벨: 일치하면 A 컬럼(원래 표기 유지), 아니면 unknown
    consensus = np.where(agree, a_raw, "unknown")
    cats = pd.unique(pd.concat([a_raw, b_raw, pd.Series(["unknown"])]).dropna())
    adata.obs["Celltype_consensus"] = pd.Categorical(consensus, categories=cats)    
    labeled_idx = np.where(adata.obs["Celltype_consensus"].astype(str) != "unknown")[0]
    scpoli_query = scPoli.load_query_data(adata=adata,reference_model=scpoli_model,labeled_indices=labeled_idx)
    scpoli_query.train(n_epochs=200,pretraining_epochs=40,eta=5)
    #adata.X = adata.X.toarray()
    #adata.X = adata.X.astype(np.float32)
    results_dict = scpoli_query.classify(adata, scale_uncertainties=True)
    annotation = results_dict['celltype_annotated']["preds"]
    uncert = results_dict['celltype_annotated']["uncert"]
    return annotation, uncert


### scPoli per sample


matching_files_filtered = pd.read_csv('/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/'+Region+'.txt',header=None)[0].to_list()

if Region == 'DLPFC':
    scpoli_model = scPoli.load("/mnt/pv_compute/dongsan/datasets/AANRI/model_scpoli/DLPFC_reference_model")
    source_adata = scpoli_model.adata.copy()
    for i in range(0,len(matching_files_filtered)): #len(matching_files_filtered)
        adata = sc.read_h5ad(matching_files_filtered[i])
        annotation, uncert = run_scpoli(scpoli_model, source_adata, adata, celltype='Celltype_MTG') # DLPFC annotation 
        adata.obs['Celltype_scpoli_sctype_consensus_MTG'] = annotation
        adata.obs['Celltype_scpoli_sctype_consensus_uncert_MTG'] = uncert
        adata.write(matching_files_filtered[i])

if Region == 'LIBD':
    scpoli_model = scPoli.load("/mnt/pv_compute/dongsan/datasets/AANRI/model_scpoli/DLPFC_reference_model")
    source_adata = scpoli_model.adata.copy()
    for i in range(0,len(matching_files_filtered)): #len(matching_files_filtered)
        adata = sc.read_h5ad(matching_files_filtered[i])
        annotation, uncert = run_scpoli(scpoli_model, source_adata, adata, celltype='Celltype_MTG') # DLPFC annotation 
        adata.obs['Celltype_scpoli_sctype_consensus_MTG'] = annotation
        adata.obs['Celltype_scpoli_sctype_consensus_uncert_MTG'] = uncert
        adata.write(matching_files_filtered[i])

if Region == 'hippo':
    scpoli_model = scPoli.load("/mnt/pv_compute/dongsan/datasets/AANRI/model_scpoli/hippo_LIBD_reference_model")
    source_adata = scpoli_model.adata.copy()
    for i in range(0,len(matching_files_filtered)): #len(matching_files_filtered)
        adata = sc.read_h5ad(matching_files_filtered[i])
        annotation, uncert = run_scpoli(scpoli_model, source_adata, adata, celltype='Celltype_hippo_LIBD') # DLPFC annotation 
        adata.obs['Celltype_scpoli_sctype_consensus_hippo_LIBD'] = annotation
        adata.obs['Celltype_scpoli_sctype_consensus_uncert_hippo_LIBD'] = uncert
        adata.write(matching_files_filtered[i])
    scpoli_model = scPoli.load("/mnt/pv_compute/dongsan/datasets/AANRI/model_scpoli/hippo_SCIENCE_reference_model")
    source_adata = scpoli_model.adata.copy()
    for i in range(0,len(matching_files_filtered)): #len(matching_files_filtered)
        adata = sc.read_h5ad(matching_files_filtered[i])
        annotation, uncert = run_scpoli(scpoli_model, source_adata, adata, celltype='Celltype_hippo_SCIENCE') # DLPFC annotation 
        adata.obs['Celltype_scpoli_sctype_consensus_hippo_SCIENCE'] = annotation
        adata.obs['Celltype_scpoli_sctype_consensus_uncert_hippo_SCIENCE'] = uncert
        adata.write(matching_files_filtered[i])

if Region == 'caudate':
    scpoli_model = scPoli.load("/mnt/pv_compute/dongsan/datasets/AANRI/model_scpoli/caudate_NCOMM_reference_model")
    source_adata = scpoli_model.adata.copy()
    for i in range(0,len(matching_files_filtered)): #len(matching_files_filtered)
        adata = sc.read_h5ad(matching_files_filtered[i])
        annotation, uncert = run_scpoli(scpoli_model, source_adata, adata, celltype='Celltype_caudate_NCOMM') # DLPFC annotation 
        adata.obs['Celltype_scpoli_sctype_consensus_caudate_NCOMM'] = annotation
        adata.obs['Celltype_scpoli_sctype_consensus_uncert_caudate_NCOMM'] = uncert
        adata.write(matching_files_filtered[i])
    scpoli_model = scPoli.load("/mnt/pv_compute/dongsan/datasets/AANRI/model_scpoli/caudate_SCIENCE_reference_model")
    source_adata = scpoli_model.adata.copy()
    for i in range(0,len(matching_files_filtered)): #len(matching_files_filtered)
        adata = sc.read_h5ad(matching_files_filtered[i])
        annotation, uncert = run_scpoli(scpoli_model, source_adata, adata, celltype='Celltype_caudate_SCIENCE') # DLPFC annotation 
        adata.obs['Celltype_scpoli_sctype_consensus_caudate_SCIENCE'] = annotation
        adata.obs['Celltype_scpoli_sctype_consensus_uncert_caudate_SCIENCE'] = uncert
        adata.write(matching_files_filtered[i])





run_harmony = '/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/harmony_integration.R'
result = subprocess.run(["Rscript", run_harmony] + [Region], capture_output=True, text=True)


# Check if the R script executed successfully
if result.returncode == 0:
    print("✅ R script executed successfully!")
else:
    print(f"❌ R script failed with return code {result.returncode}.")
    print("Error output:\n", result.stderr)





#adata_merged_harmony = sc.read('/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/'+Region+'_harmony.h5ad')
#adata_merged_harmony_top200 = sc.read('/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/'+Region+'_harmony_top200.h5ad')


matching_files_filtered = pd.read_csv('/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/'+Region+'.txt',header=None)[0].to_list()
adatas=[]
for f in matching_files_filtered:
    adata = sc.read_h5ad(f)
    adatas.append(adata)

adata_merged = sc.concat(adatas, join='outer')

adata_merged_harmony = adata_merged.copy()
if 'features' not in adata_merged_harmony.var.columns:
    adata_merged_harmony.var['features'] = adata_merged_harmony.var.index

import scvi
import scvi.model as SCVI

markers = pd.read_csv(f'/mnt/pv_compute/dongsan/datasets/AANRI/reference_for_AANRI/{Region}_marker_top_200_degs_per_celltype.csv')
markergenes = set(markers.gene)
sc.pp.highly_variable_genes(adata_merged_harmony, flavor="seurat_v3",
                            n_top_genes=2000,subset=False, batch_key='Batch')
hvg = set(adata_merged_harmony.var_names[adata_merged_harmony.var['highly_variable'].values])
#set(adata_merged_harmony.var.features.index).union(markergenes)
present_genes = set(adata_merged_harmony.var_names)
genes_to_keep_set = (hvg | markergenes) & present_genes
genes_to_keep = [g for g in adata_merged_harmony.var_names if g in genes_to_keep_set]
adata_merged_harmony_subset = adata_merged_harmony[:, genes_to_keep].copy()

sc.pp.highly_variable_genes(adata_merged_harmony, flavor="seurat_v3",
                            n_top_genes=2000,subset=False, batch_key='Batch')
hvg = set(adata_merged_harmony.var_names[adata_merged_harmony.var['highly_variable'].values])
set(adata_merged_harmony.var.features.index).union(markergenes)
present_genes = set(adata_merged_harmony.var_names)
genes_to_keep_set = (hvg | markergenes) & present_genes
genes_to_keep = [g for g in adata_merged_harmony.var_names if g in genes_to_keep_set]
adata_merged_harmony_subset = adata_merged_harmony[:, genes_to_keep].copy()    

scvi.model.SCVI.setup_anndata(adata_merged_harmony_subset,batch_key = 'Batch',
                              continuous_covariate_keys=["n_genes","pct_counts_mt","pct_counts_ribo"],
                              categorical_covariate_keys = ['Sex'])

model = scvi.model.SCVI(adata_merged_harmony_subset)
model.to_device('cuda:1')   # moves model to GPU 1
model.train(max_epochs=400,early_stopping=True)
model.adata.raw=None

adata_merged_scvi = model.adata.copy()
adata_merged_scvi.obsm["X_scVI"] = model.get_latent_representation()

sc.pp.neighbors(adata_merged_scvi, use_rep="X_scVI")
sc.tl.umap(adata_merged_scvi)
#sc.pl.umap(adata_merged_scvi)

model.save("/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/"+Region+"_scvi_top200",overwrite=True,save_anndata=True)
adata_merged_scvi.write_h5ad("/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/"+Region+"_scvi_top200/adata.h5ad")

#model = scvi.model.SCVI.load("/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/"+Region+"_scvi")
#adata_merged_scvi = model.adata.copy()



plot_umap_with_centroids(adata_merged_scvi,color='Celltype_scpoli_sctype_consensus')
plot_umap_with_centroids(adata_merged_scvi,color='Batch',annotate=False)

sc.tl.leiden(adata_merged_scvi,flavor="igraph", n_iterations=-1,directed=False,resolution=0.5)
adata_merged_scvi.obs['scvi_resolution_0.5'] = adata_merged_scvi.obs['leiden'].copy()
sc.tl.leiden(adata_merged_scvi,flavor="igraph", n_iterations=-1,directed=False,resolution=1)
adata_merged_scvi.obs['scvi_resolution_1'] = adata_merged_scvi.obs['leiden'].copy()
sc.tl.leiden(adata_merged_scvi,flavor="igraph", n_iterations=-1,directed=False,resolution=1.5)
adata_merged_scvi.obs['scvi_resolution_1.5'] = adata_merged_scvi.obs['leiden'].copy()
sc.tl.leiden(adata_merged_scvi,flavor="igraph", n_iterations=-1,directed=False,resolution=2)
adata_merged_scvi.obs['scvi_resolution_2'] = adata_merged_scvi.obs['leiden'].copy()
adata_merged_scvi.write_h5ad('/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/DLPFC_scvi/adata_with_umap.h5ad')


adata_merged_scvi=sc.read_h5ad('/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/'+Region+'_scvi/adata.h5ad')

import matplotlib
import seaborn as sns
import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
import numpy as np
import pandas as pd
import scanpy as sc
from adjustText import adjust_text  # Import adjustText for label adjustment

sns.set_style('ticks')
matplotlib.rcParams['figure.figsize'] = [4, 4]
matplotlib.rcParams['figure.dpi'] = 200

def plot_umap_with_centroids(adata, color, annotate=True, figsize=matplotlib.rcParams['figure.figsize'], dpi=matplotlib.rcParams['figure.dpi']):
    """
    Plots a UMAP projection with optional centroids annotated for the specified category.
    Ensures proper categorical coloring and prevents text overlap.
    
    Parameters:
    - adata: AnnData object containing UMAP coordinates and category annotations.
    - color: str, name of the column in `adata.obs` for annotations (e.g., 'Celltype').
    - annotate: bool, whether to annotate centroids on the plot (default: True).
    - figsize: tuple, size of the figure (default: (4, 4)).
    - dpi: int, resolution of the figure in dots per inch (default: 200).
    """
    # Convert category to categorical type if not already
    if not isinstance(adata.obs[color].dtype, pd.CategoricalDtype):
        adata.obs[color] = adata.obs[color].astype("category")
    # Generate distinct colors for all unique categories
    num_categories = adata.obs[color].nunique()
    palette = sns.color_palette("hsv", num_categories)  # Use HSV to ensure color distinction
    adata.uns[f"{color}_colors"] = [mcolors.rgb2hex(color) for color in palette]
    # Plot UMAP without legend
    sc.pl.umap(
        adata,
        color=[color],
        frameon=True,
        ncols=1,
        show=False,
        legend_loc=None,
        size=1
    )
    if annotate:
        # Calculate centroids
        umap_coordinates = adata.obsm['X_umap']
        categories = adata.obs[color]
        centroids = {cat: umap_coordinates[categories == cat].mean(axis=0) for cat in np.unique(categories)}
        # Annotate centroids on the UMAP plot with dynamic text adjustment
        ax = plt.gca()
        texts = [ax.text(centroid[0], centroid[1], cat, fontsize=10, ha='center', va='center', fontweight="bold")
                 for cat, centroid in centroids.items()]
        adjust_text(texts, ax=ax, arrowprops=dict(arrowstyle="-", lw=0.5))  # Adjust text to prevent overlap
    plt.show()


def plot_umap_side_by_side(
    adatas,
    colors,                 # 각 서브플롯에서 사용할 obs 컬럼명 리스트 (길이 = len(adatas))
    titles=None,            # 각 패널 제목
    annotate=True,
    share_palette=True,     # 같은 카테고리명이면 같은 색을 쓰고 싶다면 True
    figsize=(10, 4),
    dpi=200,
    dot_size=12,
    legend=False,           # 패널별 범례 노출 여부
):
    """
    여러 AnnData의 UMAP을 한 그림에서 좌우로 나란히 표시.

    Parameters
    ----------
    adatas : list[AnnData]
        UMAP 좌표('X_umap')와 주어진 obs 컬럼이 존재해야 함.
    colors : list[str]
        각 패널에서 색칠할 obs 컬럼명 (예: ['Celltype', 'Celltype'])
    titles : list[str] | None
        각 패널 제목
    annotate : bool
        센트로이드 텍스트 표시 여부
    share_palette : bool
        True면 전체 패널의 카테고리 이름의 합집합으로 하나의 팔레트 공유
        False면 패널별로 독립 팔레트 사용
    figsize : tuple
        전체 그림 크기
    dpi : int
        해상도
    dot_size : int
        점 크기
    legend : bool
        범례 노출 여부
    """
    assert len(adatas) == len(colors), "adatas와 colors 길이가 같아야 합니다."
    n = len(adatas)

    # 각 obs 컬럼을 카테고리로 보장
    for adata, col in zip(adatas, colors):
        if not isinstance(adata.obs[col].dtype, pd.CategoricalDtype):
            adata.obs[col] = adata.obs[col].astype("category")

    # 팔레트 준비
    if share_palette:
        # 모든 패널의 카테고리 이름(문자열) 합집합으로 팔레트 생성
        all_cats = []
        for adata, col in zip(adatas, colors):
            all_cats.extend(list(adata.obs[col].cat.categories.astype(str)))
        uniq_cats = pd.Index(pd.unique(all_cats))
        palette = sns.color_palette("hsv", len(uniq_cats))
        color_map = dict(zip(uniq_cats, [mcolors.rgb2hex(c) for c in palette]))
    else:
        color_map = None  # 패널별로 따로 생성

    fig, axes = plt.subplots(1, n, figsize=figsize, dpi=dpi, squeeze=False)
    axes = axes[0]

    for i, (adata, col) in enumerate(zip(adatas, colors)):
        ax = axes[i]

        # 패널별 팔레트
        if share_palette:
            # scanpy가 adata.uns[f"{col}_colors"]를 사용하므로 순서에 맞춰 색 리스트 세팅
            cats = adata.obs[col].cat.categories.astype(str)
            adata.uns[f"{col}_colors"] = [color_map[c] for c in cats]
        else:
            num_cats = adata.obs[col].nunique()
            pal = sns.color_palette("hsv", num_cats)
            adata.uns[f"{col}_colors"] = [mcolors.rgb2hex(c) for c in pal]

        # UMAP 그리기 (한 패널에)
        sc.pl.umap(
            adata,
            color=[col],
            frameon=True,
            ncols=1,
            show=False,
            legend_loc=None if not legend else "right margin",
            size=dot_size,
            ax=ax,
        )

        # 센트로이드 라벨
        if annotate:
            umap = adata.obsm["X_umap"]
            cats = adata.obs[col]
            # numpy 인덱싱을 위해 values로 마스크
            texts = []
            for cat in cats.cat.categories:
                mask = (cats == cat).values
                if mask.sum() == 0:
                    continue
                centroid = umap[mask].mean(axis=0)
                t = ax.text(
                    centroid[0],
                    centroid[1],
                    str(cat),
                    fontsize=9,
                    ha="center",
                    va="center",
                    fontweight="bold",
                )
                texts.append(t)
            if texts:
                adjust_text(texts, ax=ax, arrowprops=dict(arrowstyle="-", lw=0.5))

        if titles and i < len(titles):
            ax.set_title(titles[i])

    plt.tight_layout()
    return fig


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

def export_figure_to_ppt(fig, ppt_path, slide_title="UMAP comparison", margin_in=0.5, overwrite_first=False):
    """
    Save a matplotlib Figure to a new PPTX slide.
    - fig: matplotlib.figure.Figure
    - ppt_path: path to save PPTX. If exists, append a slide; else create new file.
    - slide_title: title text on the slide
    - margin_in: margin from slide edges in inches
    """
    # 1) Save the figure as a temporary PNG
    tmp_png = "/mnt/pv_compute/dongsan/tmp_plot.png"
    fig.savefig(tmp_png, dpi=300, bbox_inches="tight")
    
    if overwrite_first and os.path.exists(ppt_path):
        os.remove(ppt_path)
    
    if os.path.exists(ppt_path):
        prs = Presentation(ppt_path)
    else:
        prs = Presentation()
        prs.slide_width = Inches(10)       # ✅ Google Slides 기본 너비
        prs.slide_height = Inches(5.625)   # ✅ Google Slides 기본 높이
    # Choose a title+content layout (usually index 1)
    blank_layout = prs.slide_layouts[6]  # ✅ 완전히 빈 레이아웃
    slide = prs.slides.add_slide(blank_layout)
    
    # Add Title (textbox)
    title_box = slide.shapes.add_textbox(
        Inches(margin_in),
        Inches(margin_in),
        prs.slide_width - Inches(2*margin_in),
        Inches(0.4)
    )
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = slide_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    
    # Add the image centered with margins
    available_w = prs.slide_width - Inches(2*margin_in)
    available_h = prs.slide_height - Inches(2*margin_in) - Inches(0.7)  # account for title box
    
    pic_left = Inches(margin_in)
    pic_top = Inches(margin_in) + Inches(0.7)
    
    # Insert and let PowerPoint scale to fit the bounding box
    pic = slide.shapes.add_picture(tmp_png, pic_left, pic_top, width=available_w)
    
    # Save presentation
    prs.save(ppt_path)
    
    # Cleanup temp
    if os.path.exists(tmp_png):
        os.remove(tmp_png)


def make_global_palette(adata_ref, col, palette_name="tab20"):
    # 카테고리 보장 + 순서 고정
    if not pd.api.types.is_categorical_dtype(adata_ref.obs[col]):
        adata_ref.obs[col] = adata_ref.obs[col].astype("category")
    cats = adata_ref.obs[col].cat.categories.astype(str)
    # 팔레트 길이 확장 (필요 시 반복)
    base = sns.color_palette(palette_name, max(10, len(cats)))
    if len(cats) > len(base):
        mult = (len(cats) // len(base)) + 1
        base = (base * mult)[:len(cats)]
    colors = [mcolors.to_hex(c) for c in base[:len(cats)]]
    return list(cats), dict(zip(cats, colors))


def apply_global_palette(adata, col, global_cats, global_color_map):
    # 해당 obs 컬럼을 카테고리로 만들고, 전역 카테고리 순서로 강제
    if not pd.api.types.is_categorical_dtype(adata.obs[col]):
        adata.obs[col] = adata.obs[col].astype("category")
    adata.obs[col] = adata.obs[col].cat.set_categories(global_cats)
    # scanpy가 참조하는 팔레트 등록
    adata.uns[f"{col}_colors"] = [global_color_map.get(cat, "#808080") for cat in global_cats]

def apply_palette_to_many(adatas, col, global_cats, global_color_map):
    for A in adatas:
        if A is not None:
            apply_global_palette(A, col, global_cats, global_color_map)


Region='DLPFC'
matching_files_filtered = pd.read_csv('/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/'+Region+'.txt',header=None)[0].to_list()
adata_merged_harmony = sc.read('/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/'+Region+'_harmony.h5ad')
global_cats, global_cmap = make_global_palette(adata_merged_harmony, 'Celltype_scpoli_sctype_consensus')


for i, f in enumerate(matching_files_filtered):
    adata = sc.read_h5ad(f)
    batch_val = adata.obs.Batch.iloc[0]
    sub = adata_merged_harmony[adata_merged_harmony.obs.Batch == batch_val, :].copy()
    
    apply_palette_to_many([adata, sub], 'Celltype_scpoli_sctype_consensus', global_cats, global_cmap)
    # UMAP 두 패널 생성
    fig = plot_umap_side_by_side(
        adatas=[adata, sub],
        colors=['Celltype_scpoli_sctype_consensus', 'Celltype_scpoli_sctype_consensus'],
        titles=['No integration', 'Subplot from harmony integration'],
        annotate=True,
        share_palette=True,
        figsize=(10, 4),
        dpi=200,
        dot_size=12,
        legend=False
    )
    print([i, batch_val])
    export_figure_to_ppt(
        fig,
        "/mnt/pv_compute/dongsan/integrated.pptx",
        slide_title=str(batch_val + " (N=" + str(adata.shape[0]) + ")"),
        overwrite_first=(i == 0)  # ✅ 첫 슬라이드만 초기화
    )
    plt.close(fig)   # ✅ figure 닫기




celltype = adata_merged_harmony.obs.Celltype_scpoli_sctype_consensus.unique()

for i, f in enumerate(celltype):
    sub = adata_merged_harmony[adata_merged_harmony.obs.Celltype_scpoli_sctype_consensus ==f, :].copy()
    fig = plot_umap_side_by_side(
        adatas=[sub, adata_merged_harmony],
        colors=['Celltype_scpoli_sctype_consensus', 'Celltype_scpoli_sctype_consensus'],
        titles=['Subset', 'Subplot from harmony integration'],
        annotate=True,
        share_palette=True,
        figsize=(10, 4),
        dpi=200,
        dot_size=12,
        legend=False
        )
    export_figure_to_ppt(
        fig,
        "/mnt/pv_compute/dongsan/integrated.pptx",
        slide_title=str(f),
        overwrite_first=(i == 0)  # ✅ 첫 슬라이드만 초기화
    )
    print([i, f])
    plt.close(fig)   # ✅ figure 닫기




celltype = scpoli_model.adata.obs.celltype_annotated.unique()

for i, f in enumerate(celltype):
    sub = scpoli_model.adata[scpoli_model.adata.obs.celltype_annotated ==f, :].copy()
    fig = plot_umap_side_by_side(
        adatas=[sub, scpoli_model.adata],
        colors=['celltype_annotated', 'celltype_annotated'],
        titles=['Subset', 'Reference'],
        annotate=True,
        share_palette=True,
        figsize=(10, 4),
        dpi=200,
        dot_size=12,
        legend=False
        )
    export_figure_to_ppt(
        fig,
        "/mnt/pv_compute/dongsan/integrated.pptx",
        slide_title=str(f),
        overwrite_first=(i == 0)  # ✅ 첫 슬라이드만 초기화
    )
    print([i, f])
    plt.close(fig)   # ✅ figure 닫기






from pydrive2.auth import GoogleAuth
from pydrive2.drive import GoogleDrive
import os

# 1) OAuth 인증
gauth = GoogleAuth()
gauth.LoadClientConfigFile("/mnt/pv_compute/dongsan/client_secret_232924049563-5incek1hechaqs4ietucksh55ak1itm3.apps.googleusercontent.com.json")
gauth.settings['oauth_scope'] = ['https://www.googleapis.com/auth/drive.file']
gauth.CommandLineAuth()

# ② 구글 드라이브 객체 생성

drive = GoogleDrive(gauth)

file_path = "/mnt/pv_compute/dongsan/integrated.pptx"
folder_id = "1TUFI6FDHNwbWF4V_6GG4Sdil47KFJFyf"

meta = {
    "title": Region+'.pptx',
    "parents": [{"id": folder_id}],   # ✅ 특정 폴더에 업로드
    "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
}
f = drive.CreateFile(meta)
f.SetContentFile(file_path)
f.Upload()

print("✅ 업로드 완료!")
print("🔗 링크:", f["alternateLink"])


plot_umap_with_centroids(adata_merged_harmony,color='Celltype_scpoli_sctype_consensus')
plot_umap_with_centroids(adata_merged_harmony,color='Batch',annotate=False)






def cluster_entropy(
    adata,
    cluster_key="leiden",
    label_key="celltype_annotated",
    normalize=True,     # log(k)로 정규화 → [0,1] 범위
    base=2,             # 엔트로피 밑 (2면 bits)
    pseudocount=0.0     # 필요시 라벨 희소할 때 1e-9 같은 소량 추가 가능
):
    # 교차표 (클러스터 × 세포타입) 카운트
    ct = pd.crosstab(adata.obs[cluster_key], adata.obs[label_key])
    if pseudocount > 0:
        ct = ct + pseudocount
    # 비율
    props = ct.div(ct.sum(axis=1), axis=0).replace(0, np.nan)
    # Shannon entropy: -sum p*log(p)
    logf = (np.log2 if base == 2 else (np.log if base is np.e else lambda x: np.log(x)/np.log(base)))
    H = -(props * logf(props)).sum(axis=1)
    H = H.fillna(0.0)
    # 정규화 (최대값 log(k), k=클러스터 내 관측된 라벨 수)
    if normalize:
        k = (props > 0).sum(axis=1)
        Hmax = logf(k.replace(0, 1))  # k=0 방어
        Hn = (H / Hmax.replace(0, np.nan)).fillna(0.0)
    else:
        Hn = H
    # 지배적 세포타입과 purity
    dom_label = props.idxmax(axis=1)
    dom_prop  = props.max(axis=1)
    out = pd.DataFrame({
        "n_cells": ct.sum(axis=1),
        "n_celltypes": (props > 0).sum(axis=1),
        "entropy": H,
        "entropy_norm": Hn,
        "dominant_label": dom_label,
        "purity": dom_prop,            # 1 - impurity 느낌
    }).sort_index()
    # 원하면 비율 테이블도 같이 반환
    return out, props


def compute_entropy_for_resolutions_celltype(
    adata,
    label_key="Celltype_scpoli_sctype_consensus",
    resolutions=(0.5, 1, 1.5, 2, 2.5),
    base=2,
    normalize=True,
    pseudocount=0.0,
    key_template="harmony_snn_res.{res}",          # obs에 존재하는 클러스터 컬럼 이름 포맷
    obs_col_template="Celltype_entropy.{res}",     # 결과를 쓸 obs 컬럼 포맷
):
    ent_tables = {}      # 각 해상도별 ent_df 저장
    prop_tables = {}     # 각 해상도별 비율표 저장

    for res in resolutions:
        cluster_key = key_template.format(res=res)
        obs_col     = obs_col_template.format(res=res)

        if cluster_key not in adata.obs.columns:
            print(f"[skip] {cluster_key} not found in adata.obs")
            continue

        # 엔트로피 계산
        ent_df, prop_table = cluster_entropy(
            adata,
            cluster_key=cluster_key,
            label_key=label_key,
            normalize=normalize,
            base=base,
            pseudocount=pseudocount,
        )

        # 매핑을 위해 문자열로 통일
        ent_df.index = ent_df.index.astype(str)
        adata.obs[cluster_key] = adata.obs[cluster_key].astype(str)
        # 셀별 엔트로피를 obs에 기록
        entropy_map = ent_df["entropy_norm"].to_dict()
        adata.obs[obs_col] = adata.obs[cluster_key].map(entropy_map).astype(float)
        ent_tables[res]  = ent_df
        prop_tables[res] = prop_table
        print(f"[done] {cluster_key} → {obs_col} (clusters={ent_df.shape[0]})")
    return ent_tables, prop_tables

def compute_entropy_for_resolutions_sample(
    adata,
    label_key="Batch",
    resolutions=(0.5, 1, 1.5, 2, 2.5),
    base=2,
    normalize=True,
    pseudocount=0.0,
    key_template="harmony_snn_res.{res}",          # obs에 존재하는 클러스터 컬럼 이름 포맷
    obs_col_template="Donor_entropy.{res}",     # 결과를 쓸 obs 컬럼 포맷
):
    ent_tables = {}      # 각 해상도별 ent_df 저장
    prop_tables = {}     # 각 해상도별 비율표 저장

    for res in resolutions:
        cluster_key = key_template.format(res=res)
        obs_col     = obs_col_template.format(res=res)

        if cluster_key not in adata.obs.columns:
            print(f"[skip] {cluster_key} not found in adata.obs")
            continue

        # 엔트로피 계산
        ent_df, prop_table = cluster_entropy(
            adata,
            cluster_key=cluster_key,
            label_key=label_key,
            normalize=normalize,
            base=base,
            pseudocount=pseudocount,
        )
        # 매핑을 위해 문자열로 통일
        ent_df.index = ent_df.index.astype(str)
        adata.obs[cluster_key] = adata.obs[cluster_key].astype(str)
        # 셀별 엔트로피를 obs에 기록
        entropy_map = ent_df["entropy_norm"].to_dict()
        adata.obs[obs_col] = adata.obs[cluster_key].map(entropy_map).astype(float)
        ent_tables[res]  = ent_df
        prop_tables[res] = prop_table
        print(f"[done] {cluster_key} → {obs_col} (clusters={ent_df.shape[0]})")
    return ent_tables, prop_tables




# 실행
ent_tables, prop_tables = compute_entropy_for_resolutions_celltype(
    adata_merged_scvi,
    label_key="Celltype_scpoli_sctype_consensus_MTG",
    resolutions=(0.5, 1, 1.5, 2,2.5),
    key_template="scvi_resolution_{res}",
    normalize=False
)

# 실행
ent_tables, prop_tables = compute_entropy_for_resolutions_sample(
    adata_merged_scvi,
    label_key="Batch",
    resolutions=(0.5, 1, 1.5, 2),
    key_template="scvi_resolution_{res}",
    normalize=False
)


sc.pl.umap(
    adata_merged_scvi,
    color=["scvi_resolution_0.5", "scvi_resolution_1", "scvi_resolution_1.5", "scvi_resolution_2"],
    cmap="viridis",
    ncols=2,legend_loc=None,
)


sc.pl.umap(
    adata_merged_scvi,
    color=["Celltype_entropy.0.5", "Celltype_entropy.1", "Celltype_entropy.1.5", "Celltype_entropy.2"],
    cmap="viridis",
    ncols=2
)

sc.pl.umap(
    adata_merged_scvi,
    color=["Donor_entropy.0.5", "Donor_entropy.1", "Donor_entropy.1.5", "Donor_entropy.2"],
    cmap="viridis",
    ncols=2
)






# 실행
ent_tables, prop_tables = compute_entropy_for_resolutions_celltype(
    adata_merged_harmony,
    label_key="Celltype_scpoli_sctype_consensus_MTG",
    resolutions=(0.5, 1, 1.5, 2,2.5),
    normalize=False
)

# 실행
ent_tables, prop_tables = compute_entropy_for_resolutions_sample(
    adata_merged_harmony,
    label_key="Batch",
    resolutions=(0.5, 1, 1.5, 2,2.5),
    normalize=True
)



sc.pl.umap(
    adata_merged_harmony,
    color=["Celltype_entropy.0.5", "Celltype_entropy.1", "Celltype_entropy.1.5", "Celltype_entropy.2","Celltype_entropy.2"],
    cmap="viridis",
    ncols=2
)

sc.pl.umap(
    adata_merged_harmony,
    color=["Donor_entropy.0.5", "Donor_entropy.1", "Donor_entropy.1.5", "Donor_entropy.2","Donor_entropy.2.5"],
    cmap="viridis",
    ncols=2
)





plot_umap_with_centroids(adata_merged_harmony,color='harmony_snn_res.0.5')
plot_umap_with_centroids(adata_merged_harmony,color='Celltype_scpoli_sctype_consensus')
plot_umap_with_centroids(adata_merged_harmony,color='Batch',annotate=False)
sc.pl.umap(adata_merged_harmony,color=['Donor_entropy.0.5','Celltype_entropy.0.5'],ncols=2)
sc.pl.umap(adata_merged_harmony,color=['Celltype_scpoli_sctype_consensus_uncert','doublet_score'],ncols=2)

resolutions = [0.5, 1, 1.5, 2,2.5]
for thr in resolutions:
    col_name = f"Donor_entropy.{thr}"  # 원래 컬럼명
    new_col = f"Donor_entropy_flag_{thr}"  # 새로 만들 컬럼명
    
    adata_merged_harmony.obs[new_col] = np.where(
        (adata_merged_harmony.obs[col_name] > 0.5) ,
        True,
        False
    )

sc.pl.umap(
    adata_merged_harmony,
    color=["Donor_entropy_flag_0.5", "Donor_entropy_flag_1", "Donor_entropy_flag_1.5", 
           "Donor_entropy_flag_2","Donor_entropy_flag_2.5"],
    ncols=2
)


resolutions = [0.5, 1, 1.5, 2,2.5]
for thr in resolutions:
    col_name = f"Celltype_entropy.{thr}"  # 원래 컬럼명
    new_col = f"Celltype_entropy_flag_{thr}"  # 새로 만들 컬럼명
    
    adata_merged_harmony.obs[new_col] = np.where(
        (adata_merged_harmony.obs[col_name] < 2) ,
        True,
        False
    )

sc.pl.umap(
    adata_merged_harmony,
    color=["Celltype_entropy_flag_0.5", "Celltype_entropy_flag_1", "Celltype_entropy_flag_1.5", 
           "Celltype_entropy_flag_2","Celltype_entropy_flag_2.5"],
    ncols=2
)


if Region=='hippo':
    donor_thresh = 0.5
    ctype_thresh = 2
if Region=='DLPFC':
    donor_thresh = 0.8
    ctype_thresh = 2
if Region=='caudate':
    donor_thresh = 0.5
    ctype_thresh = 2





# 1. Donor_entropy / Celltype_entropy 각각에 대해 median 계산
donor_cols = [f"Donor_entropy.{res}" for res in resolutions]
ctype_cols = [f"Celltype_entropy.{res}" for res in resolutions]

# median 계산 후 새로운 컬럼 추가
adata_merged_harmony.obs["Donor_entropy.median"] = adata_merged_harmony.obs[donor_cols].median(axis=1)
adata_merged_harmony.obs["Celltype_entropy.median"] = adata_merged_harmony.obs[ctype_cols].median(axis=1)

# 2. median 기반으로 필터링 마스크 생성
mask = (
    (adata_merged_harmony.obs["Donor_entropy.median"] <= donor_thresh) |
    (adata_merged_harmony.obs["Celltype_entropy.median"] >= ctype_thresh)
)



adata_merged_harmony.obs["Donor_entropy_flag"] = np.where((adata_merged_harmony.obs['Donor_entropy.median'] > 0.7) ,
        True,False)
adata_merged_harmony.obs["Celltype_entropy_flag"] = np.where((adata_merged_harmony.obs['Celltype_entropy.median'] < 2.0) ,
        True,False)
sc.pl.umap(adata_merged_harmony,color="Donor_entropy_flag")
sc.pl.umap(adata_merged_harmony,color="Celltype_entropy_flag")

sns.scatterplot(x="Donor_entropy.median",y="Celltype_entropy.median",
    data=adata_merged_harmony.obs,s=20, alpha=0.6)
# 수평/수직 기준선 추가
plt.axhline(y=2.5, color='red', linestyle='--', linewidth=1)
plt.axvline(x=0.8, color='blue', linestyle='--', linewidth=1)
plt.xlabel("Donor_entropy.median")
plt.ylabel("Celltype_entropy.median")
plt.title("Entropy scatter with threshold lines")
plt.show()


# 3. 필터링된 셀 추출
filtered_cells = adata_merged_harmony[~mask].copy()

adata_merged_harmony.obs["selected"] = "Filtered out"
adata_merged_harmony.obs.loc[filtered_cells.obs_names, "selected"] = "Filtered in"

plot_umap_with_centroids(adata_merged_harmony,color='selected',annotate=False)

plot_umap_with_centroids(filtered_cells,color='Celltype_scpoli_sctype_consensus_MTG',annotate=True)


import matplotlib.pyplot as plt

# 교차표 생성
ct = pd.crosstab(
    adata_merged_harmony.obs['selected'],
    adata_merged_harmony.obs['Batch']
)

# 각 column(=Batch)별로 합을 1로 정규화
ct_norm = ct.div(ct.sum(axis=0), axis=1)

# 그리기
ax = ct_norm.T.plot(
    kind="bar",
    stacked=True,
    figsize=(16, 8),
    colormap="tab20"
)

# 그래프 꾸미기
plt.ylabel("Normalized Frequency")
plt.xlabel("Batch")
plt.title("Normalized Frequency per Batch")
plt.legend(title="Selected", bbox_to_anchor=(1.05, 1), loc="upper left")
plt.tight_layout()
plt.show()


plot_umap_with_centroids(adata_merged_harmony,color='selected',annotate=False)



plot_umap_with_centroids(filtered_cells,color='Celltype_scpoli_sctype_consensus')
sc.pl.umap(filtered_cells,color='Celltype_entropy.0.5')
sc.pl.umap(filtered_cells,color='Celltype_entropy.1')
plot_umap_with_centroids(filtered_cells,color='harmony_snn_res.0.5')
plot_umap_with_centroids(filtered_cells,color='harmony_snn_res.1')
plot_umap_with_centroids(filtered_cells,color='harmony_snn_res.1.5')
plot_umap_with_centroids(filtered_cells,color='harmony_snn_res.2')



label_key = "Celltype_scpoli_sctype_consensus_MTG"
for res in resolutions:
    cluster_key = f"harmony_snn_res.{res}"
    new_col = f"Celltype_major.{res}"
    # 1. 클러스터별 major 세포타입 계산 (filtered_cells 기준)
    ct = pd.crosstab(
        filtered_cells.obs[cluster_key],
        filtered_cells.obs[label_key]
    )
    major_labels = ct.idxmax(axis=1)
    # 2. 전체 세포에 클러스터별 major 세포타입 annotation
    filtered_cells.obs[new_col] = \
        filtered_cells.obs[cluster_key].map(major_labels)
    print(f"[res={res}] Major celltype annotation completed → {new_col}")


cols = ["Celltype_major.0.5", "Celltype_major.1", "Celltype_major.1.5", "Celltype_major.2","Celltype_major.2.5"]

# 행 단위로 가장 많이 등장한 세포타입 선택
def resolve_mode(row):
    modes = row.mode()  # 최빈값 리스트
    if len(modes) > 1:  # tie 발생 시
        return "Unknown"
    return modes.iloc[0]  # 최빈값 1개일 때

filtered_cells.obs["Celltype_major_final"] = (
    filtered_cells.obs[cols]
    .apply(resolve_mode, axis=1)
)



# 조건에 맞는 세포에 값 할당
if Region=='DLPFC':
    filtered_cells.obs["Celltype_major_final"] = (
    filtered_cells.obs["Celltype_major_final"] )


if Region=='caudate':
    filtered_cells.obs["Celltype_major_final"] = (
    filtered_cells.obs["Celltype_major_final"]
    .cat.add_categories(["VLMC"])
    )
    filtered_cells.obs.loc[
        filtered_cells.obs["harmony_snn_res.0.5"] == "14",
        "Celltype_major_final"
    ] = "VLMC"

    filtered_cells.obs.loc[
    filtered_cells.obs["Celltype_major_final"] == "Unknown",
    "Celltype_major_final"
    ] = "In:MSN_D1"


plot_umap_with_centroids(filtered_cells,color='Celltype_major_final')

order = (
    filtered_cells.obs.groupby("Celltype_scpoli_sctype_consensus_MTG")["Celltype_entropy.median"]
    .median().sort_values().index)
plt.figure(figsize=(10, 5))
sns.boxplot(
    x="Celltype_scpoli_sctype_consensus_MTG",
    y="Celltype_entropy.median",   # 혹은 y="Entropy" 등 실제 컬럼명
    data=filtered_cells.obs, palette="Set2",
    order=order
)

plt.xticks(rotation=90, ha='right')
plt.xlabel("Cell type")
plt.ylabel("Celltype Entropy (median)")
plt.title("Celltype-wise Entropy distribution")

plt.tight_layout()
plt.show()

x=pd.crosstab(filtered_cells.obs.Celltype_major_final,
              filtered_cells.obs.Batch)
sns.barplot(x.loc['Ex: L2/3 IT'])




sc.pp.normalize_total(filtered_cells, target_sum=1e4)
sc.pp.log1p(filtered_cells)


sc.pl.umap(filtered_cells, color=["PECAM1", "VWF"], use_raw=False, cmap='Reds',size=5) #Endothelial
sc.pl.umap(filtered_cells, color=["PDGFRB","RGS5","CSPG4","ACTA2","MYH11"],use_raw=False,size=5,cmap='Reds') # vascular VLMC
sc.pl.umap(filtered_cells, color=["COL1A1","COL1A2"],use_raw=False,size=5,cmap='Reds')
sc.pl.umap(filtered_cells, color=[ "COL11A1", "FAP","ACTA2","POSTN"],use_raw=False,size=5,cmap='Reds') #Fibroblast active
sc.pl.umap(filtered_cells, color=[ "DCN","PDFGRA"],use_raw=False,size=5,cmap='Reds') #Fibroblast overall




table=pd.crosstab(filtered_cells.obs.Batch,filtered_cells.obs.Celltype_major_final)
table_normalized = table.div(table.sum(axis=1), axis=0)

# Reorder the x-axis based on the values in the "Oligodendrocyte" column
table_normalized_sorted = table_normalized.sort_values(by="Oligodendrocyte", ascending=False)

# Plot a stacked bar chart
ax = table_normalized_sorted.plot(kind="bar", stacked=True, figsize=(18, 6), colormap="tab20")
# Customize plot appearance
legend = plt.legend(title="Batch", bbox_to_anchor=(1.05, 1), loc="upper left", ncol=2)

# Show the plot
plt.xlabel("Batch", fontsize=12)
plt.ylabel("Proportion", fontsize=12)
plt.title("Stacked Bar Plot of Celltype Distribution Across Batches", fontsize=14)

plt.tight_layout()
plt.show()




ax = filtered_cells.obs.Batch.value_counts()[table_normalized_sorted.index].plot(kind="bar", figsize=(18, 6), colormap="tab20")
plt.xlabel("Batch", fontsize=12)
plt.ylabel("Proportion", fontsize=12)
plt.title("Bar Plot of Nuclei Number Distribution Across Batches", fontsize=14)

plt.tight_layout()
plt.show()



age_per_batch = filtered_cells.obs.groupby("Batch")["AgeDeath"].mean()
df = table_normalized.join(age_per_batch)
corr = df.corr()

# AgeDeath와 각 cell type 비율만 추출
corr_with_age = corr["AgeDeath"].drop("AgeDeath")
print(corr_with_age)

import seaborn as sns
import matplotlib.pyplot as plt
from scipy.stats import pearsonr

fig, axes = plt.subplots(1, 2, figsize=(12,5), sharey=True)

for ax, ct, color in zip(axes, ["OPC", "Oligodendrocyte"], ["tomato", "steelblue"]):
    # 산점도 + 회귀선
    sns.regplot(x="AgeDeath", y=ct, data=df, ax=ax, color=color,
                scatter_kws={"alpha":0.7}, line_kws={"lw":2})
    
    # 상관계수
    r, p = pearsonr(df["AgeDeath"], df[ct])
    ax.set_title(f"{ct}\n r={r:.3f}, p={p:.2e}")
    ax.set_xlabel("Age at Death")
    ax.set_ylabel("Proportion")

plt.tight_layout()
plt.show()


fig, axes = plt.subplots(1,2, figsize=(12,5), sharey=True)
lowest5 = df.nsmallest(6, "AgeDeath")

for ax, ct, color in zip(axes, ["OPC", "Oligodendrocyte"], ["tomato","steelblue"]):
    sns.regplot(x="AgeDeath", y=ct, data=df, ax=ax, color=color,
                scatter_kws={"alpha":0.7}, line_kws={"lw":2})
    r, p = pearsonr(df["AgeDeath"], df[ct])
    ax.set_title(f"{ct}\n r={r:.3f}, p={p:.2e}")
    ax.set_xlabel("Age at Death")
    ax.set_ylabel("Proportion")

    texts = []
    for batch, row in lowest5.iterrows():
        texts.append(ax.text(row["AgeDeath"], row[ct], str(batch),
                             fontsize=8, color="black"))
    adjust_text(texts, ax=ax,
                arrowprops=dict(arrowstyle="->", color="gray", lw=0.5))

plt.tight_layout()
plt.show()


filtered_cells.raw = None
filtered_cells.write_h5ad(f'/mnt/pv_compute/dongsan/datasets/AANRI/integrated_filtered_final/{Region}_adata_consensus_full.h5ad')






def add_entropy_to_obs(adata, leiden_key, sample_col="Batch", prefix=None):
    """
    각 leiden cluster별 Shannon entropy를 계산하고,
    adata.obs에 entropy 값만 붙여준다.
    
    Parameters
    ----------
    adata : AnnData
    leiden_key : str
        obs에 있는 leiden cluster column 이름 (예: "harmony_snn_res.0.5")
    sample_col : str
        샘플 구분 column (예: "Batch")
    prefix : str, optional
        obs에 저장할 컬럼 접두사. 기본은 leiden_key 뒤 suffix를 따서 "entropy_rX".
    """
    # prefix 자동 생성
    if prefix is None:
        try:
            suf = str(leiden_key).split("res.")[-1]
            prefix = f"r{suf}"
        except Exception:
            prefix = leiden_key.replace(".", "_")
    out_col = f"entropy_{prefix}"
    
    # 교차표 (클러스터 × 샘플)
    clabels = adata.obs[leiden_key].astype(str)
    slabels = adata.obs[sample_col].astype(str)
    counts_df = pd.crosstab(clabels, slabels)

    # Shannon entropy per cluster
    entropy_dict = {}
    for cl, row in counts_df.iterrows():
        prop = row / row.sum()
        p = prop[prop > 0]
        H = -(p * np.log2(p)).sum()
        entropy_dict[cl] = H
    
    # 매핑해서 obs에 추가
    adata.obs[out_col] = clabels.map(entropy_dict).astype("float64")
    print(f"[{leiden_key}] → {out_col} 추가 완료.")
    return adata


resolutions = ["harmony_snn_res.0.5", "harmony_snn_res.1", "harmony_snn_res.1.5","harmony_snn_res.2"]

for r in resolutions:
    adata_merged_harmony = add_entropy_to_obs(
        adata_merged_harmony,
        leiden_key=r,
        sample_col="Batch"
    )



# resolution별 entropy 컬럼 목록
entropy_cols = ["entropy_r0.5", "entropy_r1", "entropy_r1.5", "entropy_r2"]

# 평균/최소 entropy 계산
adata_merged_harmony.obs["entropy_mean"] = adata_merged_harmony.obs[entropy_cols].mean(axis=1)
adata_merged_harmony.obs["entropy_min"]  = adata_merged_harmony.obs[entropy_cols].min(axis=1)

sc.pl.umap(adata_merged_harmony,color='entropy_mean')
sc.pl.umap(adata_merged_harmony,color='entropy_min')

adata_merged_harmony.obs["remove_flag"] = adata_merged_harmony.obs["entropy_mean"] < 4.5

print(adata_merged_harmony.obs["remove_flag"].value_counts())

# 제거 적용
adata_filtered = adata_merged_harmony[~adata_merged_harmony.obs["remove_flag"], :].copy()


sc.pl.umap(adata_merged_harmony,color='remove_flag')
plot_umap_with_centroids(adata_filtered,color="Celltype_scpoli_sctype_consensus")
sc.pl.umap(adata_filtered,color="Celltype_scpoli_sctype_consensus_uncert")




variant_df[variant_df.index==top_snp]
### differential expression
import scvi
import scvi.model as SCVI

scvi_model = scvi.model.SCVI.load('/mnt/pv_compute/dongsan/datasets/AANRI/h5ad_per_sample_filtered_FINAL_FY23_24_25/'+Region+'_scvi/')


sub = adata_merged_scvi[adata_merged_scvi.obs.Celltype_scpoli_sctype_consensus=='Oligodendrocyte'].copy()



de_change = model.differential_expression(idx1=cell_idx1, idx2=cell_idx2)















###### data analysis
import sys
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import scanpy as sc

import scvi
import os
import anndata as ad
import seaborn as sns

import warnings
import copy
import logging
import shutil
#%matplotlib inline
import scipy.sparse as sparse
from joblib import parallel_backend

import subprocess
import re

import matplotlib.colors as mcolors
from adjustText import adjust_text  # Import adjustText for label adjustment

%matplotlib inline  
sns.set_style('ticks')
matplotlib.rcParams['figure.figsize'] = [8, 8]
matplotlib.rcParams['figure.dpi'] = 200
sc.set_figure_params(figsize=(8, 8))
scvi.settings.seed = 94705
matplotlib.use("module://matplotlib_inline.backend_inline")
%config InlineBackend.print_figure_kwargs={'facecolor' : "w"}
%config InlineBackend.figure_format='retina'


def plot_umap_with_centroids(adata, color, annotate=True, figsize=matplotlib.rcParams['figure.figsize'], dpi=matplotlib.rcParams['figure.dpi']):
    """
    Plots a UMAP projection with optional centroids annotated for the specified category.
    Ensures proper categorical coloring and prevents text overlap.
    
    Parameters:
    - adata: AnnData object containing UMAP coordinates and category annotations.
    - color: str, name of the column in `adata.obs` for annotations (e.g., 'Celltype').
    - annotate: bool, whether to annotate centroids on the plot (default: True).
    - figsize: tuple, size of the figure (default: (4, 4)).
    - dpi: int, resolution of the figure in dots per inch (default: 200).
    """
    # Convert category to categorical type if not already
    if not isinstance(adata.obs[color].dtype, pd.CategoricalDtype):
        adata.obs[color] = adata.obs[color].astype("category")
    # Generate distinct colors for all unique categories
    num_categories = adata.obs[color].nunique()
    palette = sns.color_palette("hsv", num_categories)  # Use HSV to ensure color distinction
    adata.uns[f"{color}_colors"] = [mcolors.rgb2hex(color) for color in palette]
    # Plot UMAP without legend
    sc.pl.umap(
        adata,
        color=[color],
        frameon=True,
        ncols=1,
        show=False,
        legend_loc=None,
        size=1
    )
    if annotate:
        # Calculate centroids
        umap_coordinates = adata.obsm['X_umap']
        categories = adata.obs[color]
        centroids = {cat: umap_coordinates[categories == cat].mean(axis=0) for cat in np.unique(categories)}
        # Annotate centroids on the UMAP plot with dynamic text adjustment
        ax = plt.gca()
        texts = [ax.text(centroid[0], centroid[1], cat, fontsize=10, ha='center', va='center', fontweight="bold")
                 for cat, centroid in centroids.items()]
        adjust_text(texts, ax=ax, arrowprops=dict(arrowstyle="-", lw=0.5))  # Adjust text to prevent overlap
    plt.show()



def plot_category_panels(adata, category_column, basis="umap", max_cols=10, figsize_base=(5, 5), category_order=None):
    """
    Plots a figure panel across each category in the specified column, 
    showing embeddings colored by the category.
    Parameters:
    - adata: AnnData object
    - category_column: str, column in adata.obs to group by
    - basis: str, embedding basis for plotting (default: "X_mde")
    - max_cols: int, maximum number of columns in the figure grid (default: 4)
    - figsize_base: tuple, base figure size per subplot (default: (5,5))
    Returns:
    - A matplotlib figure with panels for each category
    """
    # Get unique categories and count the number of items per category
    category_counts = adata.obs[category_column].value_counts()
    category_counts = adata.obs[category_column].value_counts()
    if category_order is not None:
        unique_categories = [cat for cat in category_order if cat in category_counts.index]
    else:
        unique_categories = category_counts.index
    num_categories = len(unique_categories)
    # Determine rows and columns for subplot grid
    cols = min(max_cols, num_categories)  
    rows = (num_categories // max_cols) + (num_categories % max_cols > 0)  # Adjust rows dynamically
    # Create figure with dynamic size
    fig, axs = plt.subplots(rows, cols, figsize=(figsize_base[0] * cols, figsize_base[1] * rows))
    # Ensure axs is always a 2D array
    if rows == 1:
        axs = axs.reshape(1, -1)
    elif cols == 1:
        axs = axs.reshape(-1, 1)
    # Plot each category
    for idx, category in enumerate(unique_categories):
        ax = axs[idx // max_cols, idx % max_cols]
        count = category_counts[category]
        sc.pl.embedding(
            adata[adata.obs[category_column] == category, :], 
            basis=basis, 
            color=category_column,
            ax=ax, 
            show=False,
            title=f"{category} (n={count})",
            frameon=False
        )
    # Turn off any excess subplots
    for idx in range(num_categories, rows * cols):
        axs[idx // max_cols, idx % max_cols].axis('off')
    plt.tight_layout()
    plt.show()


Region='DLPFC'
adata = sc.read_h5ad(f'/mnt/pv_compute/dongsan/datasets/AANRI/integrated_filtered_final/{Region}.h5ad')

plot_umap_with_centroids(adata,color='AgeGroup',annotate=False)
plot_category_panels(adata,category_column='AgeGroup')
plot_umap_with_centroids(adata,color='Celltype_consensus',annotate=True)



celltype='Astrocyte'
adata_merged_raw = sc.read_h5ad(f'/mnt/pv_compute/dongsan/FOR_AANRI/subtype_analysis/{celltype}_scvi/adata_full.h5ad')


plot_category_panels(adata_merged_raw,category_column='AgeGroup')
plot_category_panels(adata_merged_raw,category_column='Region')
plot_umap_with_centroids(adata_merged_raw,color='AgeGroup',annotate=False)
plot_umap_with_centroids(adata_merged_raw,color='Region',annotate=False)

plot_category_panels(adata_merged_raw,category_column='AgeGroup',category_order=['[0-10]', '[10-20]', '[20-40]', '[40-60]', '[60- ]'])
plot_category_panels(adata_merged_raw,category_column='Region',category_order=['DLPFC','caudate','hippo'])
